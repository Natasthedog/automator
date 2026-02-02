from __future__ import annotations

# Auto-generated split from dash_app.py

import io
import base64
import logging
import copy
from difflib import SequenceMatcher
from dataclasses import dataclass, asdict, is_dataclass
from datetime import date, timedelta
from pathlib import Path
import re
import json
import pandas as pd
from openpyxl import load_workbook
from openpyxl.utils import range_boundaries, get_column_letter
import numbers
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from lxml import etree

logger = logging.getLogger(__name__)

from .pptx.charts import update_or_add_column_chart
from .pptx.text import add_table, append_paragraphs_after_label, append_text_after_label, remove_empty_placeholders, replace_text_in_slide, replace_text_in_slide_preserve_formatting, set_text_by_name
from .time_period import _format_modelling_period, _format_study_year_range, set_time_period_text
from .waterfall.inject import populate_category_waterfall
from .waterfall.targets import modelled_category_from_scope_df, target_dimensions_from_product_description

def build_pptx_from_template(
    template_bytes,
    df,
    target_brand=None,
    project_name=None,
    scope_df=None,
    product_description_df=None,
    waterfall_targets=None,
    bucket_data=None,
    modelled_in_value: str | None = None,
    metric_value: str | None = None,
):
    global _num_cache_warning_count
    _num_cache_warning_count = 0
    prs = Presentation(io.BytesIO(template_bytes))
    # Assume Slide 1 has TitleBox & SubTitle
    slide1 = prs.slides[0]
    set_text_by_name(slide1, "TitleBox", "Monthly Performance Summary")
    set_text_by_name(slide1, "SubTitle", "Auto-generated via Dash + python-pptx")
    if target_brand:
        replace_text_in_slide(slide1, "Target Brand", target_brand)
    if project_name == "MMx" and scope_df is not None:
        try:
            year_range = _format_study_year_range(scope_df)
        except Exception:
            year_range = None
        if year_range:
            replace_text_in_slide_preserve_formatting(slide1, "<RANGE>", year_range)
        generation_date = date.today().strftime("%b %d, %Y")
        replace_text_in_slide_preserve_formatting(slide1, "<DATE>", generation_date)
    remove_empty_placeholders(slide1)

    # Assume Slide 2 is for a KPI table and a chart
    slide2 = prs.slides[1] if len(prs.slides) > 1 else prs.slides.add_slide(prs.slide_layouts[5])

    # Simple KPIs (example): top 5 brands by value
    if "Brand" in df.columns and "Value" in df.columns:
        kpis = (
            df.groupby("Brand", as_index=False)["Value"].sum()
              .sort_values("Value", ascending=False)
              .head(5)
        )
        add_table(slide2, "Table_Summary", kpis)

        # Chart: share by Brand (editable)
        categories = kpis["Brand"].tolist()
        series = {"Value": kpis["Value"].tolist()}
        update_or_add_column_chart(slide2, "Chart_ShareByBrand", categories, series)

    remove_empty_placeholders(slide2)

    if project_name == "MMx" and len(prs.slides) > 3:
        slide4 = prs.slides[3]
        if scope_df is not None:
            try:
                modelled_category = modelled_category_from_scope_df(scope_df)
            except Exception:
                modelled_category = None
            if modelled_category:
                append_text_after_label(slide4, "Modelled Category:", modelled_category)
        if product_description_df is not None:
            try:
                target_dimensions = target_dimensions_from_product_description(
                    product_description_df
                )
            except Exception:
                target_dimensions = []
            if target_dimensions:
                append_paragraphs_after_label(
                    slide4,
                    "Modelled Category:",
                    target_dimensions,
                )
        if scope_df is not None:
            try:
                time_period, week_count = _format_modelling_period(df, scope_df)
            except Exception:
                time_period = None
                week_count = None
            if time_period and week_count is not None:
                set_time_period_text(slide4, "TIME PERIOD", time_period, week_count)
        remove_empty_placeholders(slide4)

    if project_name == "MMx":
        try:
            populate_category_waterfall(
                prs,
                df,
                scope_df,
                waterfall_targets,
                bucket_data,
                modelled_in_value,
                metric_value,
            )
        except Exception:
            logger.exception("Failed to populate category waterfall slides.")
            raise

    # Return bytes
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()
