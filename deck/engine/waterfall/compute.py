from __future__ import annotations

# Auto-generated split from dash_app.py

import io
import base64
import logging
import copy
from difflib import SequenceMatcher
from dataclasses import dataclass, asdict, is_dataclass
from datetime import date, timedelta
from pathlib import Path
import re
import json
import pandas as pd
from openpyxl import load_workbook
from openpyxl.utils import range_boundaries, get_column_letter
import numbers
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from lxml import etree

logger = logging.getLogger(__name__)

from ..pptx.charts import _categories_from_chart
from ..time_period import _replace_modelling_period_placeholders_in_categories
from .resolve import _normalize_text_value, _resolve_column_from_candidates, _resolve_label_from_text
from .targets import _find_column_by_candidates, _find_column_by_row_values, _target_level_labels_from_gathered_df_with_filters

DISPLAY_LABEL = {"Own": "Own", "Cross": "Competitor"}

def _normalize_target_level_labels(labels: list[str] | None) -> list[str]:
    unique_labels = []
    seen = set()
    for label in labels or []:
        if label is None:
            continue
        value = str(label).strip()
        if not value or value in seen:
            continue
        seen.add(value)
        unique_labels.append(value)
    return unique_labels



def _build_category_waterfall_df(
    gathered_df: pd.DataFrame,
    target_level_label: str | None = None,
) -> pd.DataFrame:
    header_row = gathered_df.iloc[0] if len(gathered_df) else None
    data_start_idx = 0
    vars_col, vars_idx, _ = _resolve_column_from_candidates(
        gathered_df,
        header_row,
        ["Vars", "Var", "Variable", "Variable Name", "Bucket", "Driver"],
        context="Vars/Variable",
    )
    data_start_idx = max(data_start_idx, vars_idx)
    if not vars_col:
        raise ValueError("The gatheredCN10 file is missing a Vars/Variable column for the waterfall.")

    series_candidates = {
        "Base": ["Base"],
        "Promo": ["Promo", "Promotion", "Promotions"],
        "Media": ["Media"],
        "Blanks": ["Blanks", "Blank"],
        "Positives": ["Positives", "Positive", "Pos"],
        "Negatives": ["Negatives", "Negative", "Neg"],
    }
    series_columns = {}
    missing = []
    for key, candidates in series_candidates.items():
        found, found_idx, _ = _resolve_column_from_candidates(
            gathered_df,
            header_row,
            candidates,
            context=key,
        )
        if not found:
            missing.append(key)
        else:
            series_columns[key] = found
            data_start_idx = max(data_start_idx, found_idx)
    if missing:
        raise ValueError(
            "The gatheredCN10 file is missing waterfall columns: "
            + ", ".join(missing)
        )

    label_candidates = {
        "labs-Base": ["labs-Base", "labs Base", "labels-Base", "labels Base"],
        "labs-Promo": ["labs-Promo", "labs Promo", "labels-Promo", "labels Promo"],
        "labs-Media": ["labs-Media", "labs Media", "labels-Media", "labels Media"],
        "labs-Blanks": ["labs-Blanks", "labs Blanks", "labels-Blanks", "labels Blanks"],
        "labs-Positives": [
            "labs-Positives",
            "labs Positives",
            "labels-Positives",
            "labels Positives",
        ],
        "labs-Negatives": [
            "labs-Negatives",
            "labs Negatives",
            "labels-Negatives",
            "labels Negatives",
        ],
    }
    label_columns = {}
    for key, candidates in label_candidates.items():
        found, found_idx, _ = _resolve_column_from_candidates(
            gathered_df,
            header_row,
            candidates,
            context=key,
        )
        if found:
            label_columns[key] = found
            data_start_idx = max(data_start_idx, found_idx)

    data_df = gathered_df.iloc[data_start_idx:].copy()
    if target_level_label:
        target_col, target_idx, _ = _resolve_column_from_candidates(
            gathered_df,
            header_row,
            ["Target Level Label", "Target Level"],
            context="Target Level Label",
        )
        if target_col:
            data_start_idx = max(data_start_idx, target_idx)
            data_df = gathered_df.iloc[data_start_idx:].copy()
            normalized_target = _normalize_text_value(target_level_label)
            target_series = data_df[target_col].map(_normalize_text_value)
            data_df = data_df[target_series == normalized_target]

    ordered_cols = [vars_col] + [series_columns[key] for key in series_candidates]
    ordered_cols += list(label_columns.values())
    waterfall_df = data_df.loc[:, ordered_cols].copy()
    rename_map = {vars_col: "Vars", **{v: k for k, v in series_columns.items()}}
    rename_map.update({v: k for k, v in label_columns.items()})
    waterfall_df = waterfall_df.rename(columns=rename_map)

    for key in series_candidates:
        waterfall_df[key] = pd.to_numeric(waterfall_df[key], errors="coerce").fillna(0)
    if "Negatives" in waterfall_df.columns:
        waterfall_df["Negatives"] = waterfall_df["Negatives"].abs()

    return waterfall_df


def _compute_bucket_deltas(
    data_df: pd.DataFrame,
    metadata: dict,
    bucket_config: dict[str, dict[str, list[str]]],
    year1: str,
    year2: str,
) -> list[tuple[str, float]]:
    """Compute Year2-Year1 deltas for each bucket group.

    bucket_config maps group -> {"target_labels": [...], "subheaders_included": [...]}
    """
    target_label_id = metadata.get("target_label_id")
    year_id = metadata.get("year_id")
    if not target_label_id:
        raise ValueError("The gatheredCN10 file is missing the Target Label column.")
    if not year_id:
        raise ValueError("The gatheredCN10 file is missing the Year column.")

    normalized_year1 = _normalize_text_value(year1)
    normalized_year2 = _normalize_text_value(year2)

    target_series = data_df[target_label_id].map(_normalize_text_value)
    year_series = data_df[year_id].map(_normalize_text_value)

    deltas: list[tuple[str, float]] = []
    group_order = metadata.get("group_order", [])
    ordered_groups = [group for group in group_order if group in bucket_config]
    if not ordered_groups:
        ordered_groups = list(bucket_config.keys())
    for group in ordered_groups:
        config = bucket_config.get(group, {})
        selected_cols = [
            col for col in config.get("subheaders_included", []) if col in data_df.columns
        ]
        target_labels = config.get("target_labels")
        if target_labels is None:
            target_labels = []
        if not target_labels:
            continue
        ordered_targets = []
        normalized_targets = []
        for label in target_labels:
            normalized = _normalize_text_value(label)
            if normalized and normalized not in normalized_targets:
                normalized_targets.append(normalized)
                ordered_targets.append((label, normalized))
        target_label_sequence = []
        if "own" in normalized_targets:
            target_label_sequence.append(("Own", "own"))
        if "cross" in normalized_targets:
            target_label_sequence.append(("Cross", "cross"))
        for label, normalized in ordered_targets:
            if normalized not in {"own", "cross"}:
                target_label_sequence.append((label, normalized))
        if not target_label_sequence:
            deltas.append((group, 0.0))
            continue
        if not selected_cols:
            for label, _ in target_label_sequence:
                display_label = DISPLAY_LABEL.get(label, label)
                deltas.append((f"{display_label} {group}", 0.0))
            continue
        values_df = data_df[selected_cols].apply(pd.to_numeric, errors="coerce").fillna(0)
        year1_mask = year_series == normalized_year1
        year2_mask = year_series == normalized_year2
        for label, normalized in target_label_sequence:
            target_mask = target_series == normalized
            year1_sum = values_df[year1_mask & target_mask].sum().sum()
            year2_sum = values_df[year2_mask & target_mask].sum().sum()
            display_label = DISPLAY_LABEL.get(label, label)
            deltas.append((f"{display_label} {group}", float(year2_sum - year1_sum)))
    return deltas


def _resolve_base_value_columns(gathered_df: pd.DataFrame) -> tuple[dict, int]:
    column_candidates = {
        "target_level": ["Target Level Label", "Target Level"],
        "target_label": ["Target Label", "Target", "Target Type"],
        "year": ["Year", "Model Year"],
        "actuals": ["Actuals", "Actual"],
    }
    columns = {}
    data_start_idx = 0
    header_row = gathered_df.iloc[0] if len(gathered_df) else None
    for key, candidates in column_candidates.items():
        column = _find_column_by_candidates(gathered_df, candidates)
        if not column and header_row is not None:
            column = _find_column_by_row_values(header_row, candidates)
            if column:
                data_start_idx = 1
        if not column:
            raise ValueError(
                "The gatheredCN10 file is missing the "
                f"{' / '.join(candidates)} column needed for the waterfall base."
            )
        columns[key] = column
    if columns["target_level"] == columns["target_label"]:
        alt_candidates = ["Target Level Label", "Target Level"]
        column = _find_column_by_candidates(gathered_df, alt_candidates)
        if not column and header_row is not None:
            column = _find_column_by_row_values(header_row, alt_candidates)
            if column:
                data_start_idx = 1
        if column and column != columns["target_label"]:
            columns["target_level"] = column
        else:
            raise ValueError(
                "The gatheredCN10 file needs separate Target Level Label and Target Label columns."
            )
    return columns, data_start_idx


def _waterfall_base_values(
    gathered_df: pd.DataFrame,
    target_level_label: str,
    year1: str | None = None,
    year2: str | None = None,
) -> tuple[float, float]:
    if gathered_df is None or gathered_df.empty:
        raise ValueError("The gatheredCN10 file is empty.")
    columns, data_start_idx = _resolve_base_value_columns(gathered_df)
    data_df = gathered_df.iloc[data_start_idx:]
    target_level = _normalize_text_value(target_level_label)
    target_level_series = data_df[columns["target_level"]].map(_normalize_text_value)
    target_label_series = data_df[columns["target_label"]].map(_normalize_text_value)
    year_series = data_df[columns["year"]].map(_normalize_text_value)
    actuals = pd.to_numeric(data_df[columns["actuals"]], errors="coerce").fillna(0)
    base_filter = (target_level_series == target_level) & (target_label_series == "own")
    normalized_year1 = _normalize_text_value(year1) if year1 is not None else "year1"
    normalized_year2 = _normalize_text_value(year2) if year2 is not None else "year2"
    year1_total = actuals[base_filter & (year_series == normalized_year1)].sum()
    year2_total = actuals[base_filter & (year_series == normalized_year2)].sum()
    return year1_total, year2_total


def _format_lab_base_value(value: float) -> str:
    if value is None or pd.isna(value):
        return ""
    abs_value = abs(value)
    if abs_value >= 1_000_000:
        scaled = value / 1_000_000
        suffix = "m"
    elif abs_value >= 1_000:
        scaled = value / 1_000
        suffix = "k"
    else:
        return str(int(value)) if float(value).is_integer() else str(value)
    formatted = f"{scaled:g}"
    return f"{formatted}{suffix}"


def _waterfall_base_indices(categories: list[str]) -> tuple[int, int] | None:
    earliest_idx = None
    latest_idx = None
    for idx, value in enumerate(categories):
        text = "" if value is None else str(value)
        if "<earliest date>" in text:
            earliest_idx = idx
        if "<latest date>" in text:
            latest_idx = idx
    if earliest_idx is None or latest_idx is None:
        matches = [
            idx
            for idx, value in enumerate(categories)
            if "52 w/e" in ("" if value is None else str(value)).lower()
        ]
        if len(matches) >= 2:
            earliest_idx = matches[0] if earliest_idx is None else earliest_idx
            latest_idx = matches[-1] if latest_idx is None else latest_idx
    if earliest_idx is None or latest_idx is None:
        return None
    return earliest_idx, latest_idx


def _should_update_base_series(chart_series) -> bool:
    name = getattr(chart_series, "name", "")
    if not name:
        return False
    return "base" in str(name).lower()


def _is_blanks_series(chart_series) -> bool:
    name = getattr(chart_series, "name", "")
    return "blank" in str(name).lower()


def _is_positive_series(chart_series) -> bool:
    name = getattr(chart_series, "name", "")
    return "positive" in str(name).lower()


def _is_negative_series(chart_series) -> bool:
    name = getattr(chart_series, "name", "")
    return "negative" in str(name).lower()


def _bucket_value_split(bucket_values: list[float]) -> tuple[list[float], list[float]]:
    positives: list[float] = []
    negatives: list[float] = []
    for value in bucket_values:
        if value >= 0:
            positives.append(value)
            negatives.append(0.0)
        else:
            positives.append(0.0)
            negatives.append(value)
    return positives, negatives


def _bucket_blank_values(bucket_values: list[float], base_value: float) -> list[float]:
    blanks: list[float] = []
    running_total = base_value
    for value in bucket_values:
        blanks.append(running_total)
        running_total += value
    return blanks


def _align_series_values(values: list[float], total_count: int) -> list[float]:
    if total_count <= 0:
        return values
    if len(values) < total_count:
        return values + [0.0] * (total_count - len(values))
    if len(values) > total_count:
        return values[:total_count]
    return values


def _align_label_values(values: list, total_count: int) -> list:
    if total_count <= 0:
        return values
    if len(values) < total_count:
        return values + [None] * (total_count - len(values))
    if len(values) > total_count:
        return values[:total_count]
    return values


def _sanitize_numeric_value(
    value,
    *,
    label: str | None,
    field: str,
    bucket: str | None = None,
    year1: str | None = None,
    year2: str | None = None,
) -> float:
    if value is None or pd.isna(value):
        logger.warning(
            '[waterfall][sanitize] label="%s" bucket="%s" year1="%s" year2="%s" field="%s" was=%r -> 0.0',
            label or "",
            bucket or "",
            year1 or "",
            year2 or "",
            field,
            value,
        )
        return 0.0
    if isinstance(value, str):
        try:
            return float(value)
        except ValueError as exc:
            raise ValueError(
                "Non-numeric value for "
                f'{field} (label="{label}", bucket="{bucket}", year1="{year1}", year2="{year2}"): {value!r}'
            ) from exc
    try:
        return float(value)
    except (TypeError, ValueError) as exc:
        raise ValueError(
            "Non-numeric value for "
            f'{field} (label="{label}", bucket="{bucket}", year1="{year1}", year2="{year2}"): {value!r}'
        ) from exc


def _sanitize_numeric_list(
    values: list,
    *,
    label: str | None,
    field_prefix: str,
    categories: list[str] | None = None,
    bucket_labels: list[str] | None = None,
    year1: str | None = None,
    year2: str | None = None,
) -> list[float]:
    sanitized: list[float] = []
    for idx, value in enumerate(values):
        bucket = None
        if bucket_labels and idx < len(bucket_labels):
            bucket = bucket_labels[idx]
        elif categories and idx < len(categories):
            bucket = categories[idx]
        sanitized.append(
            _sanitize_numeric_value(
                value,
                label=label,
                field=f"{field_prefix}[{idx}]",
                bucket=bucket,
                year1=year1,
                year2=year2,
            )
        )
    return sanitized


def _waterfall_series_from_gathered_df(
    gathered_df: pd.DataFrame,
    scope_df: pd.DataFrame | None,
    target_level_label: str,
) -> tuple[list[str], dict[str, list[float]], dict[str, list]] | None:
    waterfall_df = _build_category_waterfall_df(
        gathered_df,
        target_level_label=target_level_label,
    )
    if waterfall_df.empty:
        return None
    categories = (
        waterfall_df["Vars"]
        .fillna("")
        .astype(str)
        .tolist()
    )
    categories = _replace_modelling_period_placeholders_in_categories(categories, scope_df)
    series_values = {}
    for key in ["Base", "Promo", "Media", "Blanks", "Positives", "Negatives"]:
        if key in waterfall_df.columns:
            series_values[key] = (
                pd.to_numeric(waterfall_df[key], errors="coerce")
                .fillna(0)
                .astype(float)
                .tolist()
            )
    label_values: dict[str, list] = {}
    for key in [
        "labs-Base",
        "labs-Promo",
        "labs-Media",
        "labs-Blanks",
        "labs-Positives",
        "labs-Negatives",
    ]:
        if key in waterfall_df.columns:
            values = []
            for value in waterfall_df[key].tolist():
                if pd.isna(value):
                    values.append(None)
                else:
                    values.append(value)
            label_values[key] = values
    if not series_values and not label_values:
        return None
    return categories, series_values, label_values



def _apply_bucket_categories(
    categories: list[str],
    bucket_labels: list[str],
    base_indices: tuple[int, int],
) -> tuple[list[str], tuple[int, int]]:
    if not bucket_labels:
        return categories, base_indices
    start_idx, end_idx = base_indices
    insert_at = start_idx + 1
    updated = list(categories[:insert_at]) + list(bucket_labels) + list(categories[insert_at:])
    return updated, (start_idx, end_idx + len(bucket_labels))



def _apply_bucket_values(
    values: list[float],
    base_indices: tuple[int, int],
    bucket_values: list[float],
) -> list[float]:
    if not bucket_values:
        return values
    start_idx = base_indices[0]
    insert_at = start_idx + 1
    return list(values[:insert_at]) + list(bucket_values) + list(values[insert_at:])

def _apply_bucket_placeholders(
    values: list[float],
    base_indices: tuple[int, int] | None,
    bucket_count: int,
) -> list[float]:
    if bucket_count <= 0 or base_indices is None:
        return values
    start_idx = base_indices[0]
    insert_at = start_idx + 1
    return list(values[:insert_at]) + ([0.0] * bucket_count) + list(values[insert_at:])

def _build_waterfall_chart_data(
    chart,
    scope_df: pd.DataFrame | None,
    gathered_df: pd.DataFrame | None = None,
    target_level_label: str | None = None,
    bucket_labels: list[str] | None = None,
    bucket_values: list[float] | None = None,
    year1: str | None = None,
    year2: str | None = None,
) -> tuple[
    ChartData,
    list[str],
    tuple[int, int] | None,
    tuple[float, float] | None,
    list[tuple[str, list[float]]],
    dict[str, list],
]:
    gathered_override = None
    gathered_label_values: dict[str, list] = {}
    if gathered_df is not None and target_level_label:
        try:
            gathered_override = _waterfall_series_from_gathered_df(
                gathered_df,
                scope_df,
                target_level_label,
            )
        except Exception as exc:
            logger.info(
                "Skipping gatheredCN10 waterfall data for %r: %s",
                target_level_label,
                exc,
            )
            gathered_override = None
    if gathered_override is not None:
        categories, gathered_series, gathered_label_values = gathered_override
    else:
        categories = _categories_from_chart(chart)
        gathered_series = {}
        categories = _replace_modelling_period_placeholders_in_categories(categories, scope_df)
        gathered_label_values = {}
    base_indices = _waterfall_base_indices(categories)
    original_base_indices = base_indices
    bucket_labels = list(bucket_labels or [])
    bucket_values = _sanitize_numeric_list(
        list(bucket_values or []),
        label=target_level_label,
        field_prefix="bucket_values",
        bucket_labels=bucket_labels,
        year1=year1,
        year2=year2,
    )
    if bucket_labels and bucket_values:
        bucket_len = min(len(bucket_labels), len(bucket_values))
        bucket_labels = bucket_labels[:bucket_len]
        bucket_values = bucket_values[:bucket_len]
    if bucket_labels and base_indices:
        categories, base_indices = _apply_bucket_categories(
            categories,
            bucket_labels,
            base_indices,
        )
    bucket_count = len(bucket_labels)
    base_values = None
    if (
        gathered_df is not None
        and target_level_label
        and base_indices is not None
    ):
        base_values = _waterfall_base_values(
            gathered_df,
            target_level_label,
            year1=year1,
            year2=year2,
        )
        base_values = (
            _sanitize_numeric_value(
                base_values[0],
                label=target_level_label,
                field="base_values[0]",
                year1=year1,
                year2=year2,
            ),
            _sanitize_numeric_value(
                base_values[1],
                label=target_level_label,
                field="base_values[1]",
                year1=year1,
                year2=year2,
            ),
        )
    cd = ChartData()
    cd.categories = categories
    base_start_value = None
    if base_values and base_values[0] is not None:
        base_start_value = float(base_values[0])
    elif base_indices is not None:
        for series in chart.series:
            if _should_update_base_series(series):
                series_values = list(series.values)
                if base_indices[0] < len(series_values):
                    base_start_value = float(series_values[base_indices[0]])
                break
    if base_start_value is None:
        base_start_value = 0.0

    positive_bucket_values = []
    negative_bucket_values = []
    blank_bucket_values = []
    if bucket_labels and bucket_values:
        positive_bucket_values, negative_bucket_values = _bucket_value_split(bucket_values)
        blank_bucket_values = _bucket_blank_values(bucket_values, base_start_value)

    series_candidates = list(gathered_series.keys())
    series_values: list[tuple[str, list[float]]] = []
    for series in chart.series:
        values = list(series.values)
        if gathered_series:
            resolved_series = None
            try:
                resolved_series = _resolve_label_from_text(
                    str(series.name),
                    series_candidates,
                )
            except Exception as exc:
                logger.info(
                    "No gatheredCN10 series match for %r: %s",
                    series.name,
                    exc,
                )
            if resolved_series:
                values = list(gathered_series.get(resolved_series, values))
                values = _align_series_values(values, len(categories))
        if original_base_indices and bucket_labels:
            if _is_positive_series(series):
                if positive_bucket_values:
                    values = _apply_bucket_values(
                        values,
                        original_base_indices,
                        positive_bucket_values,
                    )
                else:
                    values = _apply_bucket_placeholders(
                        values,
                        original_base_indices,
                        bucket_count,
                    )
            elif _is_negative_series(series):
                if negative_bucket_values:
                    values = _apply_bucket_values(
                        values,
                        original_base_indices,
                        negative_bucket_values,
                    )
                else:
                    values = _apply_bucket_placeholders(
                        values,
                        original_base_indices,
                        bucket_count,
                    )
            elif _is_blanks_series(series):
                if blank_bucket_values:
                    values = _apply_bucket_values(
                        values,
                        original_base_indices,
                        blank_bucket_values,
                    )
                else:
                    values = _apply_bucket_placeholders(
                        values,
                        original_base_indices,
                        bucket_count,
                    )
            else:
                values = _apply_bucket_placeholders(
                    values,
                    original_base_indices,
                    bucket_count,
                )
        if base_values and base_indices:
            should_update = _should_update_base_series(series)
            if not should_update and len(chart.series) == 1:
                should_update = True
            if should_update:
                if base_indices[0] < len(values):
                    values[base_indices[0]] = base_values[0]
                if base_indices[1] < len(values):
                    values[base_indices[1]] = base_values[1]
        values = _sanitize_numeric_list(
            values,
            label=target_level_label,
            field_prefix=f"series_values[{len(series_values)}]",
            categories=categories,
            year1=year1,
            year2=year2,
        )
        cd.add_series(series.name, values)
        series_values.append((series.name, values))
    return cd, categories, base_indices, base_values, series_values, gathered_label_values


@dataclass
class WaterfallPayload:
    categories: list[str]
    series_values: list[tuple[str, list[float]]]
    base_indices: tuple[int, int] | None
    base_values: tuple[float, float] | None
    gathered_label_values: dict[str, list]


def _payload_checksum(series_values: list[tuple[str, list[float]]]) -> float:
    if not series_values:
        return 0.0
    checksum = 0.0
    if isinstance(series_values[0], tuple):
        for series_idx, (_, values) in enumerate(series_values):
            for value_idx, value in enumerate(values):
                if value is None or pd.isna(value):
                    logger.warning(
                        '[waterfall][checksum] field="series_values[%d][%d]" was=%r -> 0.0',
                        series_idx,
                        value_idx,
                        value,
                    )
                    continue
                checksum += abs(float(value))
        return checksum
    for value_idx, value in enumerate(series_values):
        if value is None or pd.isna(value):
            logger.warning(
                '[waterfall][checksum] field="values[%d]" was=%r -> 0.0',
                value_idx,
                value,
            )
            continue
        checksum += abs(float(value))
    return checksum


def _filtered_gathered_row_count(
    gathered_df: pd.DataFrame | None,
    target_level_label: str,
) -> int:
    if gathered_df is None or gathered_df.empty:
        return 0
    header_row = gathered_df.iloc[0] if len(gathered_df) else None
    data_start_idx = 0
    try:
        target_col, data_start_idx, _ = _resolve_column_from_candidates(
            gathered_df,
            header_row,
            ["Target Level Label", "Target Level"],
            context="Target Level Label",
        )
    except ValueError as exc:
        logger.info(
            "Could not resolve Target Level Label column for row counts: %s",
            exc,
        )
        target_col = None
    data_df = gathered_df.iloc[data_start_idx:]
    if not target_col:
        return len(data_df)
    normalized_target = _normalize_text_value(target_level_label)
    if not normalized_target:
        return len(data_df)
    target_series = data_df[target_col].map(_normalize_text_value)
    return int(target_series.eq(normalized_target).sum())


def _chart_data_from_payload(payload: WaterfallPayload) -> ChartData:
    cd = ChartData()
    cd.categories = payload.categories
    for name, values in payload.series_values:
        cd.add_series(name, values)
    return cd


def _derive_bucket_labels_and_values(
    gathered_df: pd.DataFrame | None,
    target_level_label: str | None,
    bucket_data: dict | None,
) -> tuple[list[str], list[float]]:
    if not bucket_data:
        return [], []
    explicit_labels = list(bucket_data.get("labels") or [])
    explicit_values = [float(value) for value in list(bucket_data.get("values") or [])]
    if explicit_labels and explicit_values:
        bucket_len = min(len(explicit_labels), len(explicit_values))
        return explicit_labels[:bucket_len], explicit_values[:bucket_len]

    bucket_config = bucket_data.get("bucket_config") or {}
    year1 = bucket_data.get("year1")
    year2 = bucket_data.get("year2")
    if gathered_df is None or gathered_df.empty or not bucket_config or not year1 or not year2:
        return [], []

    target_level_col = _find_column_by_candidates(gathered_df, ["Target Level Label", "Target Level"])
    target_label_col = _find_column_by_candidates(gathered_df, ["Target Label", "Target", "Target Type"])
    year_col = _find_column_by_candidates(gathered_df, ["Year", "Model Year"])
    if not target_level_col or not target_label_col or not year_col or not target_level_label:
        return [], []

    filtered_df = gathered_df.copy()
    normalized_target = _normalize_text_value(target_level_label)
    target_series = filtered_df[target_level_col].map(_normalize_text_value)
    filtered_df = filtered_df[target_series == normalized_target]
    if filtered_df.empty:
        return [], []

    metadata = {
        "target_label_id": target_label_col,
        "year_id": year_col,
        "group_order": list(bucket_config.keys()),
    }
    deltas = _compute_bucket_deltas(
        filtered_df,
        metadata,
        bucket_config,
        str(year1),
        str(year2),
    )
    labels = [label for label, _ in deltas]
    values = [float(value) for _, value in deltas]
    return labels, values


def compute_payload_for_label(
    gathered_df: pd.DataFrame,
    scope_df: pd.DataFrame | None,
    target_level_label: str,
    bucket_data: dict | None,
    template_chart,
) -> WaterfallPayload:
    if template_chart is None:
        raise ValueError("Template chart is required to compute waterfall payloads.")
    bucket_labels, bucket_values = _derive_bucket_labels_and_values(
        gathered_df,
        target_level_label,
        bucket_data,
    )
    (
        _,
        categories,
        base_indices,
        base_values,
        series_values,
        gathered_label_values,
    ) = _build_waterfall_chart_data(
        template_chart,
        scope_df,
        gathered_df,
        target_level_label,
        bucket_labels,
        bucket_values,
        year1=bucket_data.get("year1") if bucket_data else None,
        year2=bucket_data.get("year2") if bucket_data else None,
    )
    return WaterfallPayload(
        categories=list(categories),
        series_values=[(name, list(values)) for name, values in series_values],
        base_indices=base_indices,
        base_values=base_values,
        gathered_label_values={
            key: list(values) for key, values in gathered_label_values.items()
        },
    )


def compute_waterfall_payloads_for_all_labels(
    gathered_df: pd.DataFrame,
    scope_df: pd.DataFrame | None,
    bucket_data: dict | None,
    template_chart,
    target_labels: list[str] | None = None,
) -> dict[str, WaterfallPayload]:
    labels = _normalize_target_level_labels(target_labels)
    if not labels:
        labels = _target_level_labels_from_gathered_df_with_filters(
            gathered_df,
            year1=bucket_data.get("year1") if bucket_data else None,
            year2=bucket_data.get("year2") if bucket_data else None,
            target_labels=bucket_data.get("target_labels") if bucket_data else None,
        )
    payloads_by_label = {}
    logger.info("Precomputing waterfall payloads for %d label(s).", len(labels))
    for label in labels:
        payload = compute_payload_for_label(
            gathered_df,
            scope_df,
            label,
            bucket_data,
            template_chart,
        )
        payloads_by_label[label] = payload
        row_count = _filtered_gathered_row_count(gathered_df, label)
        checksum = _payload_checksum(payload.series_values)
        logger.info(
            "Computed waterfall payload for %r: %d categories, %d rows, checksum %.2f",
            label,
            len(payload.categories),
            row_count,
            checksum,
        )
    return payloads_by_label
