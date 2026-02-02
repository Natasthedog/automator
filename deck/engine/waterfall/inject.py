from __future__ import annotations

# Auto-generated split from dash_app.py

import io
import base64
import logging
import copy
from difflib import SequenceMatcher
from dataclasses import dataclass, asdict, is_dataclass
from datetime import date, timedelta
from pathlib import Path
import re
import json
import pandas as pd
from openpyxl import load_workbook
from openpyxl.utils import range_boundaries, get_column_letter
import numbers
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from lxml import etree

logger = logging.getLogger(__name__)

from ..pptx.chart_cache import _update_chart_label_caches, _update_waterfall_chart_caches
from ..pptx.charts import _ensure_unique_chart_parts_on_slide, _load_chart_workbook, _save_chart_workbook, _waterfall_chart_from_slide, _waterfall_chart_shape_from_slide
from ..pptx.slides import _find_slide_by_marker
from ..pptx.text import replace_placeholders_strict, replace_text_in_slide_preserve_formatting
from .compute import WaterfallPayload, _align_label_values, _build_waterfall_chart_data, _chart_data_from_payload, _format_lab_base_value, compute_waterfall_payloads_for_all_labels
from .resolve import _resolve_target_level_label_for_slide
from .targets import _normalize_column_name, _target_level_labels_from_gathered_df_with_filters

def _capture_label_columns(ws, series_names: list[str]) -> dict[int, dict[str, list]]:
    label_columns: dict[int, dict[str, list]] = {}
    series_lookup = {str(name).strip().lower() for name in series_names if name}
    for col_idx in range(2, ws.max_column + 1):
        header = ws.cell(row=1, column=col_idx).value
        if not header:
            continue
        header_text = str(header).strip().lower()
        if header_text in series_lookup:
            continue
        values = [
            ws.cell(row=row_idx, column=col_idx).value
            for row_idx in range(2, ws.max_row + 1)
        ]
        label_columns[col_idx] = {"header": header, "values": values}
    return label_columns


def _apply_label_columns(ws, label_columns: dict[int, dict[str, list]], total_rows: int) -> None:
    for col_idx, column in label_columns.items():
        ws.cell(row=1, column=col_idx, value=column["header"])
        values = column["values"]
        if len(values) < total_rows:
            values = values + [None] * (total_rows - len(values))
        for row_offset in range(total_rows):
            ws.cell(row=row_offset + 2, column=col_idx, value=values[row_offset])


def _apply_gathered_waterfall_labels(
    ws,
    label_values: dict[str, list],
    total_rows: int,
) -> set[str]:
    applied_headers: set[str] = set()
    if not label_values:
        return applied_headers
    for header, values in label_values.items():
        col_idx = _find_header_column(ws, [header])
        if col_idx is None:
            continue
        applied_headers.add(_normalize_column_name(header))
        aligned_values = _align_label_values(values, total_rows)
        for row_offset in range(total_rows):
            ws.cell(row=row_offset + 2, column=col_idx, value=aligned_values[row_offset])
    return applied_headers


def _update_lab_base_label(
    label_columns: dict[int, dict[str, list]],
    base_indices: tuple[int, int] | None,
    base_values: tuple[float, float] | None,
    total_rows: int,
) -> None:
    if base_indices is None or base_values is None:
        return
    formatted_values = [
        _format_lab_base_value(value) for value in base_values
    ]
    base_rows = list(base_indices)
    for column in label_columns.values():
        header = str(column["header"]).strip().lower()
        if header != "labs-base":
            continue
        values = column["values"]
        if len(values) < total_rows:
            values.extend([None] * (total_rows - len(values)))
        for idx, base_row in enumerate(base_rows):
            if base_row is None or base_row < 0:
                continue
            if base_row < len(values):
                values[base_row] = formatted_values[idx]
        column["values"] = values
        return


def _is_blank_cell(value) -> bool:
    if value is None:
        return True
    if isinstance(value, str) and not value.strip():
        return True
    return False


def _normalize_header_value(value: str) -> str:
    return str(value).strip().lower()


def _ensure_negatives_column_positive(ws) -> None:
    header_row = None
    header_col = None
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, max_col=ws.max_column):
        for cell in row:
            value = cell.value
            if value is None:
                continue
            if _normalize_header_value(value) == "negatives":
                header_row = cell.row
                header_col = cell.column
                break
        if header_row is not None:
            break
    if header_row is None or header_col is None:
        return

    label_col = header_col - 1 if header_col > 1 else header_col
    empty_streak = 0
    for row_idx in range(header_row + 1, ws.max_row + 1):
        label_value = ws.cell(row=row_idx, column=label_col).value
        negatives_cell = ws.cell(row=row_idx, column=header_col)
        if _is_blank_cell(label_value):
            if _is_blank_cell(negatives_cell.value):
                empty_streak += 1
                if empty_streak >= 2:
                    break
                continue
            break
        empty_streak = 0
        value = negatives_cell.value
        if isinstance(value, numbers.Number) and not isinstance(value, bool):
            negatives_cell.value = abs(value)
            continue
        if isinstance(value, str) and value.lstrip().startswith("="):
            formula = value.lstrip()[1:].strip()
            if not (formula.lower().startswith("abs(") and formula.endswith(")")):
                negatives_cell.value = f"=ABS({formula})"


def _find_header_column(ws, candidates: list[str]) -> int | None:
    normalized_columns = {}
    for col_idx in range(1, ws.max_column + 1):
        value = ws.cell(row=1, column=col_idx).value
        if value is None:
            continue
        normalized_columns[_normalize_column_name(str(value))] = col_idx
    candidate_normalized = [_normalize_column_name(candidate) for candidate in candidates]
    for candidate in candidate_normalized:
        if candidate in normalized_columns:
            return normalized_columns[candidate]
    for column_key, col_idx in normalized_columns.items():
        for candidate in candidate_normalized:
            if candidate in column_key or column_key in candidate:
                return col_idx
    from difflib import get_close_matches

    matches = get_close_matches(
        " ".join(candidate_normalized),
        list(normalized_columns.keys()),
        n=1,
        cutoff=0.75,
    )
    if matches:
        return normalized_columns[matches[0]]
    return None


def _numeric_cell_value(cell) -> float | None:
    value = cell.value
    if value is None:
        return None
    if isinstance(value, numbers.Number) and not isinstance(value, bool):
        return float(value)
    if isinstance(value, str):
        stripped = value.strip()
        if stripped.startswith("="):
            cached_value = cell.internal_value
            if isinstance(cached_value, numbers.Number) and not isinstance(cached_value, bool):
                return float(cached_value)
            cached_value = getattr(cell, "_value", None)
            if isinstance(cached_value, numbers.Number) and not isinstance(cached_value, bool):
                return float(cached_value)
            return None
        try:
            return float(stripped)
        except ValueError:
            return None
    return None


def _format_waterfall_label(value: float, sign: str) -> str:
    abs_value = abs(value)
    if abs_value >= 1_000_000:
        scaled = abs_value / 1_000_000
        suffix = "m"
    elif abs_value >= 1_000:
        scaled = abs_value / 1_000
        suffix = "k"
    else:
        scaled = abs_value
        suffix = ""
    return f"{sign}{scaled:.1f}{suffix}"


def _waterfall_labs_header_for_series(series_name: str | None) -> str | None:
    if not series_name:
        return None
    series_label = str(series_name).strip().lower()
    if "base" in series_label:
        return "labs-Base"
    if "promo" in series_label:
        return "labs-Promo"
    if "media" in series_label:
        return "labs-Media"
    if "blank" in series_label:
        return "labs-Blanks"
    if "positive" in series_label:
        return "labs-Positives"
    if "negative" in series_label:
        return "labs-Negatives"
    return None


def _resolve_waterfall_labs_column(ws, series_name: str | None) -> tuple[str | None, int | None]:
    header = _waterfall_labs_header_for_series(series_name)
    if not header:
        return None, None
    return header, _find_header_column(ws, [header])


def _update_all_waterfall_labs(
    ws,
    base_indices: tuple[int, int] | None,
    base_values: tuple[float, float] | None,
    skip_headers: set[str] | None = None,
) -> None:
    skip_headers = {value for value in (skip_headers or set()) if value}
    labs_base_col = _find_header_column(ws, ["labs-Base"])
    labs_promo_col = _find_header_column(ws, ["labs-Promo"])
    labs_media_col = _find_header_column(ws, ["labs-Media"])
    labs_blanks_col = _find_header_column(ws, ["labs-Blanks"])
    labs_positives_col = _find_header_column(ws, ["labs-Positives"])
    labs_negatives_col = _find_header_column(ws, ["labs-Negatives"])

    promo_col = _find_header_column(ws, ["Promo"])
    media_col = _find_header_column(ws, ["Media"])
    positives_col = _find_header_column(ws, ["Positives"])
    negatives_col = _find_header_column(ws, ["Negatives"])

    total_rows = ws.max_row
    if labs_base_col and _normalize_column_name("labs-Base") not in skip_headers:
        for row_idx in range(2, total_rows + 1):
            ws.cell(row=row_idx, column=labs_base_col).value = None
        if base_indices and base_values:
            formatted = [_format_lab_base_value(value) for value in base_values]
            for idx, base_row in enumerate(base_indices):
                if base_row is None or base_row < 0:
                    continue
                row_idx = base_row + 2
                if row_idx <= total_rows:
                    ws.cell(row=row_idx, column=labs_base_col).value = formatted[idx]

    for row_idx in range(2, total_rows + 1):
        if labs_promo_col and _normalize_column_name("labs-Promo") not in skip_headers:
            value = _numeric_cell_value(ws.cell(row=row_idx, column=promo_col)) if promo_col else None
            cell = ws.cell(row=row_idx, column=labs_promo_col)
            if value is None or value == 0:
                cell.value = None
            else:
                sign = "+" if value > 0 else "-"
                cell.value = _format_waterfall_label(value, sign)
        if labs_media_col and _normalize_column_name("labs-Media") not in skip_headers:
            value = _numeric_cell_value(ws.cell(row=row_idx, column=media_col)) if media_col else None
            cell = ws.cell(row=row_idx, column=labs_media_col)
            if value is None or value == 0:
                cell.value = None
            else:
                sign = "+" if value > 0 else "-"
                cell.value = _format_waterfall_label(value, sign)
        if labs_blanks_col and _normalize_column_name("labs-Blanks") not in skip_headers:
            ws.cell(row=row_idx, column=labs_blanks_col).value = None
        if labs_positives_col and _normalize_column_name("labs-Positives") not in skip_headers:
            pos_value = (
                _numeric_cell_value(ws.cell(row=row_idx, column=positives_col))
                if positives_col
                else None
            )
            cell = ws.cell(row=row_idx, column=labs_positives_col)
            if pos_value is None or pos_value == 0:
                cell.value = None
            else:
                cell.value = _format_waterfall_label(pos_value, "+")
        if labs_negatives_col and _normalize_column_name("labs-Negatives") not in skip_headers:
            neg_value = (
                _numeric_cell_value(ws.cell(row=row_idx, column=negatives_col))
                if negatives_col
                else None
            )
            cell = ws.cell(row=row_idx, column=labs_negatives_col)
            if neg_value is None or neg_value == 0:
                cell.value = None
            else:
                cell.value = _format_waterfall_label(neg_value, "-")
    logger.info(
        "Waterfall chart labels: overwrote labs columns with literal values (Base=%s, Promo=%s, Media=%s, Blanks=%s, Positives=%s, Negatives=%s).",
        bool(labs_base_col),
        bool(labs_promo_col),
        bool(labs_media_col),
        bool(labs_blanks_col),
        bool(labs_positives_col),
        bool(labs_negatives_col),
    )


def _update_waterfall_positive_negative_labels(ws) -> None:
    positives_col = _find_header_column(ws, ["Positives"])
    negatives_col = _find_header_column(ws, ["Negatives"])
    labs_positives_col = _find_header_column(ws, ["labs-Positives"])
    labs_negatives_col = _find_header_column(ws, ["labs-Negatives"])
    if not positives_col and not negatives_col:
        return
    if not labs_positives_col and not labs_negatives_col:
        return

    for row_idx in range(2, ws.max_row + 1):
        if labs_positives_col:
            pos_value = (
                _numeric_cell_value(ws.cell(row=row_idx, column=positives_col))
                if positives_col
                else None
            )
            labs_cell = ws.cell(row=row_idx, column=labs_positives_col)
            if pos_value is None or pos_value == 0:
                labs_cell.value = None
            else:
                labs_cell.value = _format_waterfall_label(pos_value, "+")
        if labs_negatives_col:
            neg_value = (
                _numeric_cell_value(ws.cell(row=row_idx, column=negatives_col))
                if negatives_col
                else None
            )
            labs_cell = ws.cell(row=row_idx, column=labs_negatives_col)
            if neg_value is None or neg_value == 0:
                labs_cell.value = None
            else:
                labs_cell.value = _format_waterfall_label(neg_value, "-")


def _format_yoy_change_text(value: float) -> str:
    if value is None or pd.isna(value):
        return "0%"
    return f"{abs(value):.0%}"


def _remove_shapes_with_text(slide, targets: list[str]) -> None:
    if not targets:
        return
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue
        text_value = shape.text_frame.text
        if any(target in text_value for target in targets):
            element = shape._element
            element.getparent().remove(element)


def _update_waterfall_yoy_arrows(
    slide,
    base_values: tuple[float, float] | None,
) -> None:
    if base_values is None:
        return
    year1_total, year2_total = base_values
    if year1_total is None or year2_total is None:
        return
    if year1_total == 0:
        pct_change = 0.0
    else:
        pct_change = (year2_total - year1_total) / year1_total
    direction = "increase" if year2_total >= year1_total else "decrease"
    remove_placeholder = "<% decrease>" if direction == "increase" else "<% increase>"
    keep_placeholder = "<% increase>" if direction == "increase" else "<% decrease>"
    _remove_shapes_with_text(slide, [remove_placeholder])
    replacement_text = f"{_format_yoy_change_text(pct_change)} {direction}"
    replace_text_in_slide_preserve_formatting(slide, keep_placeholder, replacement_text)


def _set_waterfall_chart_title(chart, label: str | None) -> None:
    if not label:
        return
    title_text = f"{label} Waterfall"
    try:
        chart.has_title = True
        chart.chart_title.text_frame.text = title_text
    except Exception:
        return


def _update_waterfall_axis_placeholders(
    prs,
    slide_selector,
    target_level_label_value: str | None,
    modelled_in_value: str | None,
    metric_value: str | None,
) -> None:
    replacements = {
        "<Target Level Label>": target_level_label_value,
        "<modelled in>": modelled_in_value,
        "<metric>": metric_value,
    }
    replace_placeholders_strict(prs, slide_selector, replacements)
    if not modelled_in_value:
        logger.warning(
            "Missing/blank value for 'Sales will be modelled in:' in Project Details."
        )
    if not metric_value:
        logger.warning(
            "Missing/blank value for 'Volume metric (unique per dataset):' in Project Details."
        )


def _add_waterfall_chart_from_template(
    slide,
    template_slide,
    scope_df: pd.DataFrame | None,
    gathered_df: pd.DataFrame | None,
    target_level_label: str | None,
    bucket_data: dict | None,
    chart_name: str,
):
    template_shape = _waterfall_chart_shape_from_slide(template_slide, chart_name)
    if template_shape is None:
        raise ValueError("Could not find the waterfall chart on the <Waterfall Template> slide.")
    template_chart = template_shape.chart
    template_series_names = [series.name for series in template_chart.series]
    label_columns = _capture_label_columns(
        _load_chart_workbook(template_chart).active,
        template_series_names,
    )
    (
        cd,
        categories,
        base_indices,
        base_values,
        _,
        gathered_label_values,
    ) = _build_waterfall_chart_data(
        template_chart,
        scope_df,
        gathered_df,
        target_level_label,
        bucket_data.get("labels") if bucket_data else None,
        bucket_data.get("values") if bucket_data else None,
        year1=bucket_data.get("year1") if bucket_data else None,
        year2=bucket_data.get("year2") if bucket_data else None,
    )
    chart_type = getattr(
        template_chart,
        "chart_type",
        getattr(XL_CHART_TYPE, "WATERFALL", XL_CHART_TYPE.COLUMN_STACKED),
    )
    chart_shape = slide.shapes.add_chart(
        chart_type,
        template_shape.left,
        template_shape.top,
        template_shape.width,
        template_shape.height,
        cd,
    )
    chart_shape.name = getattr(template_shape, "name", chart_name)
    updated_wb = _load_chart_workbook(chart_shape.chart)
    total_rows = len(categories)
    _update_lab_base_label(
        label_columns,
        base_indices,
        base_values,
        total_rows,
    )
    _apply_label_columns(updated_wb.active, label_columns, total_rows)
    _ensure_negatives_column_positive(updated_wb.active)
    applied_headers = _apply_gathered_waterfall_labels(
        updated_wb.active,
        gathered_label_values,
        total_rows,
    )
    _update_all_waterfall_labs(
        updated_wb.active,
        base_indices,
        base_values,
        skip_headers=applied_headers,
    )
    _save_chart_workbook(chart_shape.chart, updated_wb)
    _update_waterfall_chart_caches(chart_shape.chart, updated_wb, categories)
    _update_waterfall_yoy_arrows(slide, base_values)
    return chart_shape


def _waterfall_template_marker(index: int) -> str:
    if index < 0:
        raise ValueError("Waterfall template index must be non-negative.")
    if index == 0:
        return "<Waterfall Template>"
    return f"<Waterfall Template{index + 1}>"


def _available_waterfall_template_slides(prs) -> list[tuple[str, object]]:
    slides = []
    idx = 0
    while True:
        marker = _waterfall_template_marker(idx)
        slide = _find_slide_by_marker(prs, marker)
        if slide is None:
            break
        slides.append((marker, slide))
        idx += 1
    return slides


def _normalize_target_level_labels(labels: list[str] | None) -> list[str]:
    unique_labels = []
    seen = set()
    for label in labels or []:
        if label is None:
            continue
        value = str(label).strip()
        if not value or value in seen:
            continue
        seen.add(value)
        unique_labels.append(value)
    return unique_labels


def _set_waterfall_slide_header(slide, label: str, marker_text: str | None = None) -> None:
    title_text = label
    replaced = False
    if marker_text:
        replaced = replace_text_in_slide_preserve_formatting(slide, marker_text, title_text)
        marker_plain = marker_text.strip("<>")
        replaced = (
            replace_text_in_slide_preserve_formatting(slide, marker_plain, title_text)
            or replaced
        )
    replaced = replace_text_in_slide_preserve_formatting(
        slide, "<Waterfall Template>", title_text
    ) or replaced
    replaced = (
        replace_text_in_slide_preserve_formatting(slide, "Waterfall Template", title_text)
        or replaced
    )
    if replaced:
        return
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.text = title_text
            return


def _update_waterfall_chart(
    slide,
    payload: WaterfallPayload,
) -> None:
    chart_shapes = [shape for shape in slide.shapes if shape.has_chart]
    if not chart_shapes:
        raise ValueError("Could not find the waterfall chart on the <Waterfall Template> slide.")
    for chart_shape in chart_shapes:
        chart = chart_shape.chart
        series_names = [series.name for series in chart.series]
        label_columns = _capture_label_columns(_load_chart_workbook(chart).active, series_names)
        cd = _chart_data_from_payload(payload)
        chart.replace_data(cd)
        updated_wb = _load_chart_workbook(chart)
        total_rows = len(payload.categories)
        _update_lab_base_label(
            label_columns,
            payload.base_indices,
            payload.base_values,
            total_rows,
        )
        _apply_label_columns(updated_wb.active, label_columns, total_rows)
        _ensure_negatives_column_positive(updated_wb.active)
        applied_headers = _apply_gathered_waterfall_labels(
            updated_wb.active,
            payload.gathered_label_values,
            total_rows,
        )
        _update_all_waterfall_labs(
            updated_wb.active,
            payload.base_indices,
            payload.base_values,
            skip_headers=applied_headers,
        )
        _save_chart_workbook(chart, updated_wb)
        chart_type = getattr(
            chart,
            "chart_type",
            getattr(XL_CHART_TYPE, "WATERFALL", XL_CHART_TYPE.COLUMN_STACKED),
        )
        if chart_type == getattr(XL_CHART_TYPE, "WATERFALL", XL_CHART_TYPE.COLUMN_STACKED):
            _update_waterfall_chart_caches(chart, updated_wb, payload.categories)
        else:
            _update_chart_label_caches(chart, updated_wb)
    _update_waterfall_yoy_arrows(slide, payload.base_values)


def _resolve_target_level_label_value(
    gathered_df: pd.DataFrame | None,
    waterfall_targets: list[str] | None,
    bucket_data: dict | None,
) -> str | None:
    selected = [label for label in (waterfall_targets or []) if label and str(label).strip()]
    if selected:
        if len(selected) > 1:
            joined = ", ".join(str(label).strip() for label in selected)
            logger.info(
                "Multiple Target Level Label values selected; using joined string: %s",
                joined,
            )
            return joined
        return str(selected[0]).strip()
    if gathered_df is None or gathered_df.empty:
        return None
    year1 = bucket_data.get("year1") if bucket_data else None
    year2 = bucket_data.get("year2") if bucket_data else None
    target_labels = bucket_data.get("target_labels") if bucket_data else None
    labels = _target_level_labels_from_gathered_df_with_filters(
        gathered_df,
        year1=year1,
        year2=year2,
        target_labels=target_labels,
    )
    if not labels:
        return None
    if len(labels) > 1:
        joined = ", ".join(labels)
        logger.info(
            "Multiple Target Level Label values derived from gatheredCN10; using joined string: %s",
            joined,
        )
        return joined
    return labels[0]


def populate_category_waterfall(
    prs,
    gathered_df: pd.DataFrame,
    scope_df: pd.DataFrame | None = None,
    target_labels: list[str] | None = None,
    bucket_data: dict | None = None,
    modelled_in_value: str | None = None,
    metric_value: str | None = None,
):
    labels = _normalize_target_level_labels(target_labels)
    if not labels:
        labels = _target_level_labels_from_gathered_df_with_filters(
            gathered_df,
            year1=bucket_data.get("year1") if bucket_data else None,
            year2=bucket_data.get("year2") if bucket_data else None,
            target_labels=bucket_data.get("target_labels") if bucket_data else None,
        )
    if not labels:
        return

    available_slides = _available_waterfall_template_slides(prs)
    available_count = len(available_slides)
    if available_count == 0:
        raise ValueError("Could not find the <Waterfall Template> slide in the template.")
    if len(labels) > available_count:
        raise ValueError(
            "Need {needed} waterfall slides but only found {available} "
            "(<Waterfall Template>...<Waterfall Template{available}>) in template. "
            "Please add more duplicated template slides or use a larger template deck.".format(
                needed=len(labels),
                available=available_count,
            )
        )

    template_chart = _waterfall_chart_from_slide(available_slides[0][1], "Waterfall Template")
    if template_chart is None:
        raise ValueError("Could not find the waterfall chart on the <Waterfall Template> slide.")
    payloads_by_label = compute_waterfall_payloads_for_all_labels(
        gathered_df,
        scope_df,
        bucket_data,
        template_chart,
        target_labels=labels,
    )

    seen_partnames: set[str] = set()
    remaining_labels = labels.copy()
    for idx in range(len(labels)):
        marker_text, slide = available_slides[idx]
        resolved_label = _resolve_target_level_label_for_slide(slide, remaining_labels)
        if resolved_label is None:
            if not remaining_labels:
                raise ValueError("No remaining Target Level Label values to assign.")
            resolved_label = remaining_labels[0]
            logger.info(
                "No slide title/name found for %s; using ordered Target Level Label %r.",
                marker_text,
                resolved_label,
            )
        if resolved_label in remaining_labels:
            remaining_labels.remove(resolved_label)
        if resolved_label not in payloads_by_label:
            raise ValueError(
                f"Missing precomputed waterfall payload for Target Level Label {resolved_label!r}."
            )
        _ensure_unique_chart_parts_on_slide(slide, seen_partnames)
        _update_waterfall_axis_placeholders(
            prs,
            slide,
            target_level_label_value=resolved_label,
            modelled_in_value=modelled_in_value,
            metric_value=metric_value,
        )
        _update_waterfall_chart(
            slide,
            payloads_by_label[resolved_label],
        )
        _set_waterfall_slide_header(slide, resolved_label, marker_text=marker_text)
