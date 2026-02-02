from __future__ import annotations

# Auto-generated split from dash_app.py

import io
import base64
import logging
import copy
from difflib import SequenceMatcher
from dataclasses import dataclass, asdict, is_dataclass
from datetime import date, timedelta
from pathlib import Path
import re
import json
import pandas as pd
from openpyxl import load_workbook
from openpyxl.utils import range_boundaries, get_column_letter
import numbers
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from lxml import etree

logger = logging.getLogger(__name__)

from .pptx.text import _normalize_label

class CompanyWeekMapper:
    """
    Maps between a sequential internal 'company_week' integer and YEARWK (YYYYWW),
    assuming YEARWK is an ISO week (ISO-8601).

    Provide at least one anchor mapping:
        (anchor_company_week -> anchor_yearwk)

    Optionally provide a second anchor to validate the mapping.
    """
    anchor_company_week: int
    anchor_yearwk: int
    check_company_week: int | None = None
    check_yearwk: int | None = None

    @staticmethod
    def _yearwk_to_monday(yearwk: int) -> date:
        y, w = divmod(yearwk, 100)
        if not (1 <= w <= 53):
            raise ValueError(f"Invalid YEARWK week number: {yearwk}")
        # Monday of ISO week
        return date.fromisocalendar(y, w, 1)

    @staticmethod
    def _monday_to_yearwk(d: date) -> int:
        iso_y, iso_w, _ = d.isocalendar()
        return iso_y * 100 + iso_w

    def __post_init__(self):
        # Optional consistency check with second anchor
        if (self.check_company_week is None) ^ (self.check_yearwk is None):
            raise ValueError("Provide both check_company_week and check_yearwk, or neither.")

        if self.check_company_week is not None:
            a_date = self._yearwk_to_monday(self.anchor_yearwk)
            delta = self.check_company_week - self.anchor_company_week
            derived = self._monday_to_yearwk(a_date + timedelta(weeks=delta))
            if derived != self.check_yearwk:
                raise ValueError(
                    f"Anchors inconsistent: derived {derived} but expected {self.check_yearwk}."
                )

    def to_yearwk(self, company_week: int) -> int:
        a_date = self._yearwk_to_monday(self.anchor_yearwk)
        delta_weeks = company_week - self.anchor_company_week
        out_date = a_date + timedelta(weeks=delta_weeks)
        return self._monday_to_yearwk(out_date)

    def to_company_week(self, yearwk: int) -> int:
        a_date = self._yearwk_to_monday(self.anchor_yearwk)
        target_date = self._yearwk_to_monday(yearwk)
        delta_days = (target_date - a_date).days

        if delta_days % 7 != 0:
            # Should never happen because both are Mondays, but keep it safe
            raise ValueError("Non-week-aligned difference; check inputs.")
        delta_weeks = delta_days // 7
        return self.anchor_company_week + delta_weeks


def _find_company_week_value(scope_df: pd.DataFrame, label: str):
    if scope_df is None or scope_df.empty or scope_df.shape[1] < 2:
        raise ValueError("Scope file must include labels in column A and company weeks in column B.")
    label_normalized = _normalize_label(label)
    candidates = []
    for _, row in scope_df.iterrows():
        cell_value = row.iloc[0]
        if pd.isna(cell_value):
            continue
        cell_label = _normalize_label(str(cell_value))
        candidates.append((cell_label, row))
        if label_normalized in cell_label or cell_label in label_normalized:
            value = row.iloc[1]
            if pd.isna(value):
                raise ValueError(f"Missing company week value for '{label}'.")
            return value
    if candidates:
        from difflib import get_close_matches

        labels = [candidate_label for candidate_label, _ in candidates]
        matches = get_close_matches(label_normalized, labels, n=1, cutoff=0.7)
        if matches:
            matched_label = matches[0]
            for candidate_label, row in candidates:
                if candidate_label == matched_label:
                    value = row.iloc[1]
                    if pd.isna(value):
                        raise ValueError(f"Missing company week value for '{label}'.")
                    return value
    raise ValueError(f"Could not find '{label}' in the scope file.")


def _coerce_yearwk(value) -> int:
    if pd.isna(value):
        raise ValueError("Missing company week value for modelling period.")
    raw = str(value).strip()
    if not raw:
        raise ValueError("Missing company week value for modelling period.")
    try:
        numeric_value = int(float(raw))
    except ValueError:
        digits = "".join(ch for ch in raw if ch.isdigit())
        if not digits:
            raise ValueError("Company week value must be a YYYYWW-style week number or a company week.")
        numeric_value = int(digits)

    if len(str(numeric_value)) <= 4:
        mapper = CompanyWeekMapper(
            anchor_company_week=2455,
            anchor_yearwk=202638,
            check_company_week=2470,
            check_yearwk=202653,
        )
        return mapper.to_yearwk(numeric_value)

    year, week = divmod(numeric_value, 100)
    if year <= 0 or not (1 <= week <= 53):
        raise ValueError("Company week value must be a YYYYWW-style week number or a company week.")
    return numeric_value


def _format_modelling_period(data_df: pd.DataFrame, scope_df: pd.DataFrame) -> tuple[str, int]:
    start_company_week = _find_company_week_value(scope_df, "First week of modelling")
    end_company_week = _find_company_week_value(scope_df, "Last week of modelling")
    start_yearwk = _coerce_yearwk(start_company_week)
    end_yearwk = _coerce_yearwk(end_company_week)
    start_date = CompanyWeekMapper._yearwk_to_monday(start_yearwk)
    end_date = CompanyWeekMapper._yearwk_to_monday(end_yearwk) + timedelta(days=6)
    week_count = ((end_date - start_date).days // 7) + 1
    return f"{start_date:%b %d, %Y} - {end_date:%b %d, %Y}", week_count


def _format_study_year_range(scope_df: pd.DataFrame) -> str:
    start_company_week = _find_company_week_value(scope_df, "First week of modelling")
    end_company_week = _find_company_week_value(scope_df, "Last week of modelling")
    start_yearwk = _coerce_yearwk(start_company_week)
    end_yearwk = _coerce_yearwk(end_company_week)
    start_year = CompanyWeekMapper._yearwk_to_monday(start_yearwk).year
    end_year = CompanyWeekMapper._yearwk_to_monday(end_yearwk).year
    if start_year == end_year:
        return str(start_year)
    return f"{start_year}-{end_year}"


def set_time_period_text(slide, label_text, time_period, week_count):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        if label_text not in text_frame.text:
            continue
        label_index = None
        for idx, paragraph in enumerate(text_frame.paragraphs):
            if label_text in paragraph.text:
                label_index = idx
                break
        if label_index is None:
            continue
        paragraphs = list(text_frame.paragraphs)
        if label_index + 1 < len(paragraphs):
            value_paragraph = paragraphs[label_index + 1]
        else:
            value_paragraph = text_frame.add_paragraph()
        value_paragraph.text = f"{time_period} (number of weeks = {week_count})"
        for extra_paragraph in paragraphs[label_index + 2:]:
            extra_paragraph.text = ""
        return True
    return False


def _modelling_period_bounds(scope_df: pd.DataFrame) -> tuple[date, date]:
    start_company_week = _find_company_week_value(scope_df, "First week of modelling")
    end_company_week = _find_company_week_value(scope_df, "Last week of modelling")
    start_yearwk = _coerce_yearwk(start_company_week)
    end_yearwk = _coerce_yearwk(end_company_week)
    earliest_yearwk = min(start_yearwk, end_yearwk)
    latest_yearwk = max(start_yearwk, end_yearwk)
    earliest_date = CompanyWeekMapper._yearwk_to_monday(earliest_yearwk)
    latest_date = CompanyWeekMapper._yearwk_to_monday(latest_yearwk) + timedelta(days=6)
    return earliest_date, latest_date


def _replace_modelling_period_placeholders_in_categories(
    categories: list[str],
    scope_df: pd.DataFrame | None,
) -> list[str]:
    if scope_df is None or not categories:
        return categories
    try:
        start_date, end_date = _modelling_period_bounds(scope_df)
    except Exception:
        return categories
    earliest = start_date.strftime("%b %d, %Y")
    latest = end_date.strftime("%b %d, %Y")
    updated = []
    for value in categories:
        text = "" if value is None else str(value)
        if "<earliest date>" in text or "<latest date>" in text:
            text = text.replace("<earliest date>", earliest)
            text = text.replace("<latest date>", latest)
        updated.append(text)
    return updated
