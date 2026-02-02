from __future__ import annotations

from pathlib import Path

import pandas as pd
from django.core.management.base import BaseCommand
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from deck.engine.waterfall.compute import compute_waterfall_payloads_for_all_labels
from deck.engine.waterfall.payloads import _waterfall_payloads_to_json
from deck.engine.waterfall.targets import _target_level_labels_from_gathered_df_with_filters
from deck.engine.waterfall.validate import validate_payloads_or_raise


def _load_dataframe(path: Path) -> pd.DataFrame:
    if path.suffix.lower() in {".xlsx", ".xls", ".xlsb"}:
        read_options = {}
        if path.suffix.lower() == ".xlsb":
            read_options["engine"] = "pyxlsb"
        return pd.read_excel(path, **read_options)
    if path.suffix.lower() == ".csv":
        return pd.read_csv(path)
    raise ValueError(f"Unsupported file format for {path}.")


def _load_scope_dataframe(path: Path) -> pd.DataFrame | None:
    if path.suffix.lower() not in {".xlsx", ".xlsb"}:
        raise ValueError("Scope file must be .xlsx or .xlsb.")
    read_options = {"sheet_name": "Overall Information"}
    if path.suffix.lower() == ".xlsb":
        read_options["engine"] = "pyxlsb"
    scope_df = pd.read_excel(path, **read_options)
    return None if scope_df.empty else scope_df


def _build_minimal_template_chart():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = ChartData()
    chart_data.categories = ["Placeholder"]
    for name in ["Base", "Promo", "Media", "Blanks", "Positives", "Negatives"]:
        chart_data.add_series(name, (0,))
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(3),
        chart_data,
    )
    return chart_shape.chart


class Command(BaseCommand):
    help = "Dump computed waterfall payloads to JSON for debugging."

    def add_arguments(self, parser):
        parser.add_argument(
            "--gathered",
            required=True,
            help="Path to gatheredCN10 CSV/XLSX/XLSB file.",
        )
        parser.add_argument(
            "--scope",
            required=False,
            help="Optional path to scope workbook (Overall Information sheet).",
        )
        parser.add_argument(
            "--out",
            required=True,
            help="Path to write JSON payloads.",
        )

    def handle(self, *args, **options):
        gathered_path = Path(options["gathered"]).expanduser()
        scope_path = options.get("scope")
        out_path = Path(options["out"]).expanduser()

        gathered_df = _load_dataframe(gathered_path)
        scope_df = _load_scope_dataframe(Path(scope_path)) if scope_path else None

        labels = _target_level_labels_from_gathered_df_with_filters(gathered_df)
        if not labels:
            raise ValueError("No Target Level Label values found in gathered data.")

        template_chart = _build_minimal_template_chart()
        payloads_by_label = compute_waterfall_payloads_for_all_labels(
            gathered_df,
            scope_df,
            bucket_data=None,
            template_chart=template_chart,
            target_labels=labels,
        )
        validate_payloads_or_raise(payloads_by_label, labels)

        out_path.write_text(
            _waterfall_payloads_to_json(payloads_by_label),
            encoding="utf-8",
        )
        self.stdout.write(self.style.SUCCESS(f"Wrote waterfall payloads to {out_path}"))
