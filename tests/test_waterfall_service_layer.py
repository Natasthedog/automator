from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pandas as pd
from lxml import etree
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from dataclasses import dataclass
from deck_automation.services.waterfall_service_layer import (
    ArtifactStore,
    WaterfallChartUpdater,
    WaterfallGenerationRequest,
    WaterfallOrchestrator,
    WaterfallPayloadBuilder,
    WaterfallPlaceholderService,
    WaterfallSlideMapper,
)


@dataclass
class _Payload:
    categories: list[str]
    series_values: list[tuple[str, list[float]]]
    base_indices: tuple[int, int] | None
    base_values: tuple[float, float] | None
    gathered_label_values: dict[str, list]


def _payload() -> _Payload:
    series_names = ["Base", "Promo", "Media", "Blanks", "Positives", "Negatives"]
    return _Payload(
        categories=["Base 2022", "Base 2023"],
        series_values=[(name, [10.0, 12.0]) for name in series_names],
        base_indices=(0, 1),
        base_values=(100.0, 120.0),
        gathered_label_values={"labs-Base": ["100", "120"]},
    )


def _build_template(path: Path, labels: list[str] | None = None) -> Path:
    prs = Presentation()
    layout = prs.slide_layouts[6]
    chart_data = ChartData()
    chart_data.categories = ["Base 2022", "Base 2023"]
    for name in ["Base", "Promo", "Media", "Blanks", "Positives", "Negatives"]:
        chart_data.add_series(name, (1, 2))
    labels = labels or ["Alpha", "Beta"]
    for idx, label in enumerate(labels):
        marker = "<Waterfall Template>" if idx == 0 else f"<Waterfall Template{idx + 1}>"
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(6), Inches(0.5))
        title.text_frame.text = marker
        slide.name = label
        placeholder = slide.shapes.add_textbox(Inches(0.3), Inches(0.7), Inches(7), Inches(0.5))
        placeholder.text_frame.text = "<Target Level Label> <modelled in> <metric>"
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED,
            Inches(1),
            Inches(1.3),
            Inches(6),
            Inches(3),
            chart_data,
        )
        chart_shape.name = "Waterfall Template"
    prs.save(path)
    return path


def test_payload_builder_default_labels_when_empty_selection():
    gathered = pd.DataFrame(
        [
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year1"},
            {"Target Level Label": "Beta", "Target Label": "Cross", "Year": "Year2"},
        ]
    )
    labels = WaterfallPayloadBuilder().resolve_target_labels(
        gathered,
        target_labels=[],
        bucket_data={"year1": "Year1", "year2": "Year2", "target_labels": ["Own", "Cross"]},
    )
    assert labels == ["Alpha", "Beta"]


def test_template_slide_mapping_order_and_resolution(tmp_path: Path):
    template_path = _build_template(tmp_path / "mapping.pptx", labels=["Beta", "Alpha"])
    prs = Presentation(template_path)
    mapper = WaterfallSlideMapper()
    slides = mapper.find_available_template_slides(prs)
    assert [marker for marker, _ in slides] == ["<Waterfall Template>", "<Waterfall Template2>"]
    remaining = ["Beta", "Alpha"]
    resolved = mapper.resolve_label_for_slide(slides[0][1], remaining, slides[0][0])
    assert resolved == "Beta"


def test_placeholder_service_replaces_markers(tmp_path: Path):
    template_path = _build_template(tmp_path / "placeholders.pptx")
    prs = Presentation(template_path)
    slide = prs.slides[0]
    service = WaterfallPlaceholderService()
    service.update_axis_placeholders(prs, slide, "Alpha", "Volume", "Units")
    service.set_slide_header(slide, "Alpha", marker_text="<Waterfall Template>")
    text_values = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    assert any("Alpha Volume Units" in text for text in text_values)
    assert any(text.strip() == "Alpha" for text in text_values)


def test_chart_replace_data_tripwire_calls_cache_refresh(tmp_path: Path):
    template_path = _build_template(tmp_path / "tripwire.pptx")
    prs = Presentation(template_path)
    slide = prs.slides[0]
    chart = next(shape.chart for shape in slide.shapes if shape.has_chart)
    updater = WaterfallChartUpdater()
    calls: list[str] = []

    original_replace = chart.replace_data

    def wrapped_replace(data):
        calls.append("replace")
        return original_replace(data)

    chart.replace_data = wrapped_replace  # type: ignore[method-assign]

    original_load = updater._load_chart_workbook
    original_save = updater._save_chart_workbook
    original_cache = updater._update_chart_label_caches
    original_waterfall_cache = updater._update_waterfall_chart_caches

    def wrapped_load(target_chart):
        calls.append("load")
        return original_load(target_chart)

    def wrapped_save(target_chart, workbook):
        calls.append("save")
        return original_save(target_chart, workbook)

    def wrapped_cache(target_chart, workbook):
        calls.append("cache")
        return original_cache(target_chart, workbook)

    def wrapped_waterfall_cache(target_chart, workbook, categories):
        calls.append("cache")
        return original_waterfall_cache(target_chart, workbook, categories)

    updater._load_chart_workbook = wrapped_load  # type: ignore[method-assign]
    updater._save_chart_workbook = wrapped_save  # type: ignore[method-assign]
    updater._update_chart_label_caches = wrapped_cache  # type: ignore[method-assign]
    updater._update_waterfall_chart_caches = wrapped_waterfall_cache  # type: ignore[method-assign]

    updater.update_slide_charts(slide, _payload())

    assert "replace" in calls and "save" in calls and "cache" in calls
    assert calls.index("replace") < calls.index("save") < calls.index("cache")


def test_golden_cache_xml_contains_str_cache_and_dlbl_range_cache(tmp_path: Path):
    template_path = _build_template(tmp_path / "cache.pptx")
    prs = Presentation(template_path)
    slide = prs.slides[0]
    chart = next(shape.chart for shape in slide.shapes if shape.has_chart)
    root = chart.part._element
    ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    ns_c15 = "http://schemas.microsoft.com/office/drawing/2012/chart"
    ns = {"c": ns_c}
    series = root.find(".//c:ser", namespaces=ns)
    dlbls = etree.SubElement(series, f"{{{ns_c}}}dLbls")
    dlbl_range = etree.SubElement(dlbls, f"{{{ns_c15}}}datalabelsRange")
    etree.SubElement(dlbl_range, f"{{{ns_c15}}}f").text = "Sheet1!$B$2:$B$3"

    WaterfallChartUpdater().update_slide_charts(slide, _payload())

    xml = etree.tostring(chart.part._element)
    assert b"strCache" in xml
    assert b"dlblRangeCache" in xml


def test_orchestrator_generates_valid_pptx_and_relationships(tmp_path: Path):
    template_path = _build_template(tmp_path / "orchestrator.pptx", labels=["Alpha", "Beta"])
    gathered_df = pd.read_csv(Path("tests/fixtures/gathered_min.csv"))
    store = ArtifactStore(base_tmp_dir=tmp_path / "tmp", durable_root=tmp_path / "durable")
    orchestrator = WaterfallOrchestrator(artifact_store=store)
    result = orchestrator.generate(
        WaterfallGenerationRequest(
            template_path=template_path,
            gathered_df=gathered_df,
            target_labels=[],
            bucket_data={"year1": "Year1", "year2": "Year2", "target_labels": ["Own"]},
            modelled_in_value="Volume",
            metric_value="Units",
            job_id="job-1",
        )
    )
    assert result.selected_labels == ["Alpha", "Beta"]
    assert result.temp_output_path.exists()
    with zipfile.ZipFile(result.temp_output_path, "r") as archive:
        names = set(archive.namelist())
    assert "[Content_Types].xml" in names
    assert "ppt/slides/_rels/slide1.xml.rels" in names
