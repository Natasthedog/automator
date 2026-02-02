from __future__ import annotations

import zipfile
from pathlib import Path

import pandas as pd
from lxml import etree
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches

from deck.engine.waterfall.inject import populate_category_waterfall
from deck.engine.pptx.charts import _waterfall_chart_from_slide


def _build_shared_chart_template(tmp_path: Path) -> Path:
    prs = Presentation()
    layout = prs.slide_layouts[6]
    chart_data = ChartData()
    chart_data.categories = ["Base 2022", "Base 2023"]
    for name in ["Base", "Promo", "Media", "Blanks", "Positives", "Negatives"]:
        chart_data.add_series(name, (1, 2))

    for idx, marker in enumerate(["<Waterfall Template>", "<Waterfall Template2>"]):
        slide = prs.slides.add_slide(layout)
        textbox = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.2),
            Inches(5),
            Inches(0.5),
        )
        textbox.text = marker
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(3),
            chart_data,
        )
        chart_shape.name = "Waterfall Template"

    template_path = tmp_path / "shared_chart_template.pptx"
    prs.save(template_path)

    with zipfile.ZipFile(template_path, "r") as zipf:
        contents = {name: zipf.read(name) for name in zipf.namelist()}

    rels_ns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
    slide1_rels = etree.fromstring(contents["ppt/slides/_rels/slide1.xml.rels"])
    slide2_rels = etree.fromstring(contents["ppt/slides/_rels/slide2.xml.rels"])
    chart_target = None
    for rel in slide1_rels.xpath("//r:Relationship", namespaces=rels_ns):
        if rel.get("Type") == RT.CHART:
            chart_target = rel.get("Target")
            break
    if not chart_target:
        raise AssertionError("Could not locate chart relationship on slide1.")
    for rel in slide2_rels.xpath("//r:Relationship", namespaces=rels_ns):
        if rel.get("Type") == RT.CHART:
            rel.set("Target", chart_target)

    contents["ppt/slides/_rels/slide2.xml.rels"] = etree.tostring(
        slide2_rels,
        xml_declaration=True,
        encoding="UTF-8",
        standalone="yes",
    )

    with zipfile.ZipFile(template_path, "w") as zipf:
        for name, data in contents.items():
            zipf.writestr(name, data)

    return template_path


def test_injection_creates_unique_chart_parts(tmp_path: Path):
    template_path = _build_shared_chart_template(tmp_path)
    prs = Presentation(template_path)

    chart1 = _waterfall_chart_from_slide(prs.slides[0], "Waterfall Template")
    chart2 = _waterfall_chart_from_slide(prs.slides[1], "Waterfall Template")
    assert str(chart1.part.partname) == str(chart2.part.partname)

    gathered_df = pd.read_csv(Path("tests/fixtures/gathered_min.csv"))
    populate_category_waterfall(
        prs,
        gathered_df,
        target_labels=["Alpha", "Beta"],
        bucket_data=None,
    )

    chart1 = _waterfall_chart_from_slide(prs.slides[0], "Waterfall Template")
    chart2 = _waterfall_chart_from_slide(prs.slides[1], "Waterfall Template")
    assert str(chart1.part.partname) != str(chart2.part.partname)
    assert (
        str(chart1.part.chart_workbook.xlsx_part.partname)
        != str(chart2.part.chart_workbook.xlsx_part.partname)
    )
