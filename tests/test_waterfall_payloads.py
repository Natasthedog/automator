from __future__ import annotations

import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from deck.engine.waterfall.compute import (
    _payload_checksum,
    compute_waterfall_payloads_for_all_labels,
)


def _build_template_chart():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = ChartData()
    chart_data.categories = ["<earliest date>", "Bridge", "<latest date>"]
    for name in ["Base", "Promo", "Media", "Blanks", "Positives", "Negatives"]:
        chart_data.add_series(name, (0,))
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(3),
        chart_data,
    )
    return chart_shape.chart


def test_payloads_differ_across_labels():
    gathered_df = pd.DataFrame(
        [
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year1", "Actuals": 100},
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year2", "Actuals": 130},
            {"Target Level Label": "Beta", "Target Label": "Own", "Year": "Year1", "Actuals": 300},
            {"Target Level Label": "Beta", "Target Label": "Own", "Year": "Year2", "Actuals": 360},
        ]
    )
    template_chart = _build_template_chart()

    payloads = compute_waterfall_payloads_for_all_labels(
        gathered_df,
        scope_df=None,
        bucket_data=None,
        template_chart=template_chart,
    )

    assert set(payloads.keys()) == {"Alpha", "Beta"}

    checksums = [_payload_checksum(payload.series_values) for payload in payloads.values()]
    assert len(set(checksums)) > 1

    for payload in payloads.values():
        assert payload.categories
        assert payload.series_values
        assert all(values for _, values in payload.series_values)


def test_payloads_use_template_categories_without_gathered_vars_columns():
    gathered_df = pd.DataFrame(
        [
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year1", "Actuals": 100},
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year2", "Actuals": 130},
            {"Target Level Label": "Beta", "Target Label": "Own", "Year": "Year1", "Actuals": 90},
            {"Target Level Label": "Beta", "Target Label": "Own", "Year": "Year2", "Actuals": 120},
        ]
    )
    template_chart = _build_template_chart()

    payloads = compute_waterfall_payloads_for_all_labels(
        gathered_df,
        scope_df=None,
        bucket_data=None,
        template_chart=template_chart,
    )

    assert payloads["Alpha"].categories[0] != ""
    assert payloads["Alpha"].base_values == (100.0, 130.0)
    assert payloads["Beta"].base_values == (90.0, 120.0)


def test_payload_checksums_differ_for_multiple_target_level_labels_without_series_columns():
    gathered_df = pd.DataFrame(
        [
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year1", "Actuals": 100},
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year2", "Actuals": 130},
            {"Target Level Label": "Beta", "Target Label": "Own", "Year": "Year1", "Actuals": 300},
            {"Target Level Label": "Beta", "Target Label": "Own", "Year": "Year2", "Actuals": 360},
        ]
    )
    template_chart = _build_template_chart()

    payloads = compute_waterfall_payloads_for_all_labels(
        gathered_df,
        scope_df=None,
        bucket_data=None,
        template_chart=template_chart,
    )

    checksums = {label: _payload_checksum(payload.series_values) for label, payload in payloads.items()}
    assert checksums["Alpha"] != checksums["Beta"]


def test_bucket_config_computes_distinct_values_per_target_level_label():
    gathered_df = pd.DataFrame(
        [
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year1", "Actuals": 100, "Price": 10},
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year2", "Actuals": 130, "Price": 25},
            {"Target Level Label": "Alpha", "Target Label": "Cross", "Year": "Year1", "Actuals": 50, "Price": 5},
            {"Target Level Label": "Alpha", "Target Label": "Cross", "Year": "Year2", "Actuals": 70, "Price": 15},
            {"Target Level Label": "Beta", "Target Label": "Own", "Year": "Year1", "Actuals": 80, "Price": 2},
            {"Target Level Label": "Beta", "Target Label": "Own", "Year": "Year2", "Actuals": 120, "Price": 6},
            {"Target Level Label": "Beta", "Target Label": "Cross", "Year": "Year1", "Actuals": 45, "Price": 1},
            {"Target Level Label": "Beta", "Target Label": "Cross", "Year": "Year2", "Actuals": 50, "Price": 3},
        ]
    )
    template_chart = _build_template_chart()

    payloads = compute_waterfall_payloads_for_all_labels(
        gathered_df,
        scope_df=None,
        bucket_data={
            "year1": "Year1",
            "year2": "Year2",
            "bucket_config": {
                "Price": {
                    "target_labels": ["Own", "Cross"],
                    "subheaders_included": ["Price"],
                }
            },
        },
        template_chart=template_chart,
    )

    assert "Own Price" in payloads["Alpha"].categories
    assert "Own Price" in payloads["Beta"].categories
    alpha_checksum = _payload_checksum(payloads["Alpha"].series_values)
    beta_checksum = _payload_checksum(payloads["Beta"].series_values)
    assert alpha_checksum != beta_checksum


def test_explicit_bucket_values_are_fallback_when_config_cannot_be_resolved():
    gathered_df = pd.DataFrame(
        [
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year1", "Actuals": 100},
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year2", "Actuals": 130},
        ]
    )
    template_chart = _build_template_chart()

    payloads = compute_waterfall_payloads_for_all_labels(
        gathered_df,
        scope_df=None,
        bucket_data={
            "year1": "Year1",
            "year2": "Year2",
            "labels": ["Own Price"],
            "values": [7.0],
            "bucket_config": {
                "Price": {
                    "target_labels": ["Own"],
                    "subheaders_included": ["Price"],
                }
            },
        },
        template_chart=template_chart,
    )

    assert "Own Price" in payloads["Alpha"].categories
