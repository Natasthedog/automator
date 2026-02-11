from __future__ import annotations

import io
import logging
import shutil
import zipfile
from dataclasses import dataclass
from pathlib import Path
from uuid import uuid4

import pandas as pd
from django.conf import settings
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

from deck.engine.pptx.chart_cache import _update_chart_label_caches, _update_waterfall_chart_caches
from deck.engine.pptx.charts import _ensure_unique_chart_parts_on_slide, _load_chart_workbook, _save_chart_workbook
from deck.engine.waterfall.compute import WaterfallPayload, _chart_data_from_payload, compute_waterfall_payloads_for_all_labels
from deck.engine.waterfall.inject import (
    _apply_gathered_waterfall_labels,
    _apply_label_columns,
    _available_waterfall_template_slides,
    _capture_label_columns,
    _ensure_negatives_column_positive,
    _set_waterfall_slide_header,
    _update_all_waterfall_labs,
    _update_lab_base_label,
    _update_waterfall_axis_placeholders,
    _update_waterfall_yoy_arrows,
)
from deck.engine.waterfall.resolve import _resolve_target_level_label_for_slide
from deck.engine.waterfall.targets import _target_level_labels_from_gathered_df_with_filters

logger = logging.getLogger(__name__)


@dataclass
class WaterfallGenerationRequest:
    template_path: Path
    gathered_df: pd.DataFrame
    scope_df: pd.DataFrame | None = None
    target_labels: list[str] | None = None
    bucket_data: dict | None = None
    modelled_in_value: str | None = None
    metric_value: str | None = None
    job_id: str | None = None


@dataclass
class WaterfallGenerationResult:
    job_id: str
    temp_output_path: Path
    storage_key: str
    selected_labels: list[str]


class WaterfallPayloadBuilder:
    def resolve_target_labels(
        self,
        gathered_df: pd.DataFrame,
        target_labels: list[str] | None,
        bucket_data: dict | None,
    ) -> list[str]:
        labels = [str(value).strip() for value in (target_labels or []) if str(value).strip()]
        if labels:
            return list(dict.fromkeys(labels))
        derived = _target_level_labels_from_gathered_df_with_filters(
            gathered_df,
            year1=bucket_data.get("year1") if bucket_data else None,
            year2=bucket_data.get("year2") if bucket_data else None,
            target_labels=bucket_data.get("target_labels") if bucket_data else None,
        )
        if not derived:
            derived = _target_level_labels_from_gathered_df_with_filters(gathered_df)
        logger.info("Resolved %d target labels from gatheredCN10 defaults.", len(derived))
        return derived

    def build_payload_for_label(
        self,
        gathered_df: pd.DataFrame,
        scope_df: pd.DataFrame | None,
        bucket_data: dict | None,
        template_chart,
        label: str,
    ) -> WaterfallPayload:
        payloads = compute_waterfall_payloads_for_all_labels(
            gathered_df,
            scope_df,
            bucket_data,
            template_chart,
            target_labels=[label],
        )
        return payloads[label]

    def build_payloads_for_all_labels(
        self,
        gathered_df: pd.DataFrame,
        scope_df: pd.DataFrame | None,
        bucket_data: dict | None,
        template_chart,
        target_labels: list[str] | None,
    ) -> dict[str, WaterfallPayload]:
        resolved = self.resolve_target_labels(gathered_df, target_labels, bucket_data)
        return compute_waterfall_payloads_for_all_labels(
            gathered_df,
            scope_df,
            bucket_data,
            template_chart,
            target_labels=resolved,
        )


class WaterfallSlideMapper:
    def find_available_template_slides(self, prs) -> list[tuple[str, object]]:
        return _available_waterfall_template_slides(prs)

    def resolve_label_for_slide(self, slide, remaining_labels: list[str], marker_text: str) -> str:
        resolved_label = _resolve_target_level_label_for_slide(slide, remaining_labels)
        if resolved_label is None:
            if not remaining_labels:
                raise ValueError("No remaining Target Level Label values to assign.")
            resolved_label = remaining_labels[0]
            logger.info(
                "No slide title/name found for %s; using ordered Target Level Label %r.",
                marker_text,
                resolved_label,
            )
        return resolved_label

    def ensure_slide_capacity(self, available_count: int, requested_count: int) -> None:
        if available_count == 0:
            raise ValueError("Could not find the <Waterfall Template> slide in the template.")
        if requested_count > available_count:
            raise ValueError(
                "Need {needed} waterfall slides but only found {available} "
                "(<Waterfall Template>...<Waterfall Template{available}>) in template. "
                "Please add more duplicated template slides or use a larger template deck.".format(
                    needed=requested_count,
                    available=available_count,
                )
            )


class WaterfallChartUpdater:
    def update_slide_charts(self, slide, payload: WaterfallPayload) -> None:
        chart_shapes = [shape for shape in slide.shapes if shape.has_chart]
        if not chart_shapes:
            raise ValueError("Could not find the waterfall chart on the <Waterfall Template> slide.")
        for chart_shape in chart_shapes:
            chart = chart_shape.chart
            series_names = [series.name for series in chart.series]
            label_columns = _capture_label_columns(self._load_chart_workbook(chart).active, series_names)
            chart.replace_data(_chart_data_from_payload(payload))
            workbook = self._load_chart_workbook(chart)
            total_rows = len(payload.categories)
            _update_lab_base_label(label_columns, payload.base_indices, payload.base_values, total_rows)
            _apply_label_columns(workbook.active, label_columns, total_rows)
            _ensure_negatives_column_positive(workbook.active)
            applied_headers = _apply_gathered_waterfall_labels(
                workbook.active,
                payload.gathered_label_values,
                total_rows,
            )
            _update_all_waterfall_labs(
                workbook.active,
                payload.base_indices,
                payload.base_values,
                skip_headers=applied_headers,
            )
            self._save_chart_workbook(chart, workbook)
            chart_type = getattr(chart, "chart_type", getattr(XL_CHART_TYPE, "WATERFALL", XL_CHART_TYPE.COLUMN_STACKED))
            if chart_type == getattr(XL_CHART_TYPE, "WATERFALL", XL_CHART_TYPE.COLUMN_STACKED):
                self._update_waterfall_chart_caches(chart, workbook, payload.categories)
            else:
                self._update_chart_label_caches(chart, workbook)
        _update_waterfall_yoy_arrows(slide, payload.base_values)

    def _load_chart_workbook(self, chart):
        return _load_chart_workbook(chart)

    def _save_chart_workbook(self, chart, workbook) -> None:
        _save_chart_workbook(chart, workbook)

    def _update_waterfall_chart_caches(self, chart, workbook, categories) -> None:
        logger.info("Refreshing waterfall cache for %d categories.", len(categories))
        _update_waterfall_chart_caches(chart, workbook, categories)

    def _update_chart_label_caches(self, chart, workbook) -> None:
        logger.info("Refreshing standard chart label cache.")
        _update_chart_label_caches(chart, workbook)


class WaterfallPlaceholderService:
    def update_axis_placeholders(
        self,
        prs,
        slide,
        target_level_label_value: str,
        modelled_in_value: str | None,
        metric_value: str | None,
    ) -> None:
        _update_waterfall_axis_placeholders(
            prs,
            slide,
            target_level_label_value=target_level_label_value,
            modelled_in_value=modelled_in_value,
            metric_value=metric_value,
        )

    def set_slide_header(self, slide, resolved_label: str, marker_text: str | None = None) -> None:
        _set_waterfall_slide_header(slide, resolved_label, marker_text=marker_text)


class ArtifactStore:
    def __init__(self, base_tmp_dir: Path | None = None, durable_root: Path | None = None) -> None:
        self.base_tmp_dir = base_tmp_dir or Path("/tmp/decks")
        self.durable_root = durable_root or (Path(settings.MEDIA_ROOT) / "decks")

    def write_temp_artifact(self, job_id: str, filename: str, content: bytes) -> Path:
        output_dir = self.base_tmp_dir / str(job_id)
        output_dir.mkdir(parents=True, exist_ok=True)
        output_path = output_dir / filename
        output_path.write_bytes(content)
        return output_path

    def promote_output(self, job_id: str, source_path: Path) -> str:
        self.durable_root.mkdir(parents=True, exist_ok=True)
        destination = self.durable_root / f"{job_id}-{source_path.name}"
        shutil.copyfile(source_path, destination)
        return str(destination.relative_to(self.durable_root.parent))


class WaterfallOrchestrator:
    def __init__(
        self,
        payload_builder: WaterfallPayloadBuilder | None = None,
        slide_mapper: WaterfallSlideMapper | None = None,
        chart_updater: WaterfallChartUpdater | None = None,
        placeholder_service: WaterfallPlaceholderService | None = None,
        artifact_store: ArtifactStore | None = None,
    ) -> None:
        self.payload_builder = payload_builder or WaterfallPayloadBuilder()
        self.slide_mapper = slide_mapper or WaterfallSlideMapper()
        self.chart_updater = chart_updater or WaterfallChartUpdater()
        self.placeholder_service = placeholder_service or WaterfallPlaceholderService()
        self.artifact_store = artifact_store or ArtifactStore()

    def generate(self, request: WaterfallGenerationRequest) -> WaterfallGenerationResult:
        job_id = request.job_id or str(uuid4())
        prs = Presentation(request.template_path)
        labels = self.payload_builder.resolve_target_labels(
            request.gathered_df,
            request.target_labels,
            request.bucket_data,
        )
        available_slides = self.slide_mapper.find_available_template_slides(prs)
        self.slide_mapper.ensure_slide_capacity(len(available_slides), len(labels))
        template_chart = next((shape.chart for shape in available_slides[0][1].shapes if shape.has_chart), None)
        if template_chart is None:
            raise ValueError("Could not find the waterfall chart on the <Waterfall Template> slide.")
        payloads = self.payload_builder.build_payloads_for_all_labels(
            request.gathered_df,
            request.scope_df,
            request.bucket_data,
            template_chart,
            labels,
        )
        seen_partnames: set[str] = set()
        remaining_labels = labels.copy()
        for idx, (marker_text, slide) in enumerate(available_slides[: len(labels)]):
            resolved_label = self.slide_mapper.resolve_label_for_slide(slide, remaining_labels, marker_text)
            if resolved_label in remaining_labels:
                remaining_labels.remove(resolved_label)
            _ensure_unique_chart_parts_on_slide(slide, seen_partnames)
            self.placeholder_service.update_axis_placeholders(
                prs,
                slide,
                resolved_label,
                request.modelled_in_value,
                request.metric_value,
            )
            self.chart_updater.update_slide_charts(slide, payloads[resolved_label])
            self.placeholder_service.set_slide_header(slide, resolved_label, marker_text=marker_text)
            logger.info("Updated waterfall slide %d (%s) with label %s", idx + 1, marker_text, resolved_label)

        stream = io.BytesIO()
        prs.save(stream)
        pptx_bytes = stream.getvalue()
        with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r"):
            pass
        temp_path = self.artifact_store.write_temp_artifact(job_id, "waterfall_output.pptx", pptx_bytes)
        storage_key = self.artifact_store.promote_output(job_id, temp_path)
        return WaterfallGenerationResult(
            job_id=job_id,
            temp_output_path=temp_path,
            storage_key=storage_key,
            selected_labels=labels,
        )
