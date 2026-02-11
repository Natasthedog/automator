from __future__ import annotations

import io
import json

import pandas as pd
from django.core.files.uploadedfile import SimpleUploadedFile
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from django.test import TestCase
from django.urls import reverse

from deck_automation.services.readers import read_df
from deck_automation.services.waterfall_payloads import compute_waterfall_payloads_for_all_labels


class DeckAutomationViewsTests(TestCase):
    def _minimal_gathered_df(self) -> pd.DataFrame:
        return pd.DataFrame(
            [
                {
                    "Target Level Label": "Alpha",
                    "Target Label": "Own",
                    "Year": "Year1",
                    "Actuals": 100,
                    "Vars": 1,
                    "Base": 10,
                    "Promo": 2,
                    "Media": 3,
                    "Blanks": 4,
                    "Positives": 5,
                    "Negatives": -6,
                },
                {
                    "Target Level Label": "Alpha",
                    "Target Label": "Own",
                    "Year": "Year2",
                    "Actuals": 110,
                    "Vars": 1,
                    "Base": 11,
                    "Promo": 2,
                    "Media": 3,
                    "Blanks": 4,
                    "Positives": 5,
                    "Negatives": -6,
                },
                {
                    "Target Level Label": "Beta",
                    "Target Label": "Own",
                    "Year": "Year1",
                    "Actuals": 200,
                    "Vars": 1,
                    "Base": 20,
                    "Promo": 4,
                    "Media": 6,
                    "Blanks": 8,
                    "Positives": 10,
                    "Negatives": -12,
                },
                {
                    "Target Level Label": "Beta",
                    "Target Label": "Own",
                    "Year": "Year2",
                    "Actuals": 210,
                    "Vars": 1,
                    "Base": 21,
                    "Promo": 4,
                    "Media": 6,
                    "Blanks": 8,
                    "Positives": 10,
                    "Negatives": -12,
                },
            ]
        )

    def _csv_upload(self) -> SimpleUploadedFile:
        csv_content = self._minimal_gathered_df().to_csv(index=False).encode("utf-8")
        return SimpleUploadedFile("gathered.csv", csv_content, content_type="text/csv")

    def _xlsx_upload(self) -> SimpleUploadedFile:
        buffer = io.BytesIO()
        self._minimal_gathered_df().to_excel(buffer, index=False)
        return SimpleUploadedFile(
            "gathered.xlsx",
            buffer.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    def _template_chart(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        chart_data = ChartData()
        chart_data.categories = ["Placeholder"]
        for name in ["Base", "Promo", "Media", "Blanks", "Positives", "Negatives"]:
            chart_data.add_series(name, (0,))
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED,
            Inches(1),
            Inches(1),
            Inches(6),
            Inches(3),
            chart_data,
        )
        return chart_shape.chart

    def test_read_df_handles_csv_and_xlsx(self):
        csv_df = read_df(self._csv_upload())
        xlsx_df = read_df(self._xlsx_upload())

        self.assertIn("Target Level Label", csv_df.columns)
        self.assertIn("Target Level Label", xlsx_df.columns)
        self.assertEqual(len(csv_df), 4)
        self.assertEqual(len(xlsx_df), 4)

    def test_payloads_accept_alias_headers_and_embedded_header_row(self):
        gathered_df = pd.DataFrame(
            [
                {
                    "Target Level Label": "Alpha",
                    "Target Label": "Own",
                    "Year": "Year1",
                    "Actuals": 100,
                },
                {
                    "Target Level Label": "Alpha",
                    "Target Label": "Own",
                    "Year": "Year2",
                    "Actuals": 110,
                },
            ]
        )

        payloads = compute_waterfall_payloads_for_all_labels(
            gathered_df,
            scope_df=None,
            bucket_data=None,
            template_chart=self._template_chart(),
        )

        self.assertIn("Alpha", payloads)
        payload = payloads["Alpha"]
        self.assertEqual(payload.categories, ["Placeholder"])
        self.assertEqual(payload.base_values, (100.0, 110.0))

    def test_read_df_detects_header_row_on_row_2_for_excel(self):
        raw_df = pd.DataFrame(
            [
                ["ignore", "ignore", "ignore", "ignore"],
                ["Target Level Label", "Target Label", "Year", "Actuals"],
                ["Alpha", "Own", "Year1", 100],
                ["Alpha", "Own", "Year2", 120],
            ]
        )
        buffer = io.BytesIO()
        raw_df.to_excel(buffer, index=False, header=False)
        upload = SimpleUploadedFile(
            "row2.xlsx",
            buffer.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

        gathered_df = read_df(upload)

        self.assertEqual(gathered_df.attrs.get("detected_header_row_index"), 1)
        self.assertEqual(
            list(gathered_df.columns),
            ["Target Level Label", "Target Label", "Year", "Actuals"],
        )
        self.assertEqual(gathered_df.iloc[0]["Target Level Label"], "Alpha")

    def test_deck_automation_page_loads(self):
        response = self.client.get(reverse("deck-automation"))

        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "Deck Automation (MVP)")

    def test_root_redirects_to_deck_automation(self):
        response = self.client.get("/")

        self.assertEqual(response.status_code, 302)
        self.assertEqual(response.headers.get("Location"), "/deck-automation/")

    def test_post_computes_payloads_and_renders_summary(self):
        response = self.client.post(
            reverse("deck-automation"),
            data={"gathered_cn10": self._csv_upload(), "template_choice": "MMx"},
        )

        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "Computed")
        self.assertContains(response, "Download payloads JSON")


    def test_post_defaults_template_choice_to_mmx(self):
        response = self.client.post(
            reverse("deck-automation"),
            data={"gathered_cn10": self._csv_upload()},
        )

        self.assertEqual(response.status_code, 200)
        self.assertContains(response, "Computed")

    def test_download_endpoint_returns_json_payloads(self):
        post_response = self.client.post(
            reverse("deck-automation"),
            data={"gathered_cn10": self._csv_upload(), "template_choice": "MMM"},
        )
        self.assertEqual(post_response.status_code, 200)

        download_id = post_response.context["download_id"]
        download_response = self.client.get(
            reverse("deck-automation-download", kwargs={"download_id": download_id})
        )

        self.assertEqual(download_response.status_code, 200)
        self.assertEqual(download_response["Content-Type"], "application/json")

        payload = json.loads(download_response.content.decode("utf-8"))
        self.assertIn("Alpha", payload)
        self.assertIn("Beta", payload)
