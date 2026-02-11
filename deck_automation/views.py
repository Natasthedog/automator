from __future__ import annotations

import json
import logging
from pathlib import Path
from uuid import uuid4

from django.conf import settings
from django.http import Http404, HttpResponse
from django.shortcuts import render
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from deck.engine.pptx.charts import _waterfall_chart_from_slide
from deck.engine.waterfall.inject import _available_waterfall_template_slides

from deck_automation.services.readers import read_df, read_scope_df
from deck_automation.services.waterfall_payloads import (
    compute_waterfall_payloads_for_all_labels,
    payload_checksum,
    waterfall_payloads_to_json,
)

logger = logging.getLogger(__name__)

SESSION_PAYLOADS_KEY = "deck_automation_payloads"
WATERFALL_TEMPLATE_MARKER = "<Waterfall Template>"
TEMPLATE_OPTIONS = {
    "MMx": "MMx.pptx",
    "MMM": "MMM.pptx",
    "PnP": "PnP.pptx",
}


def _find_template_chart(template_source):
    if template_source is None:
        return None
    prs = Presentation(template_source)
    available_slides = _available_waterfall_template_slides(prs)
    if available_slides:
        chart = _waterfall_chart_from_slide(available_slides[0][1], "Waterfall Template")
        if chart is not None:
            return chart
    logger.info("Falling back to first chart found in template deck.")
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_chart:
                return shape.chart
    return None



def _default_template_chart():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = ChartData()
    chart_data.categories = ["<earliest date>", "Bridge", "<latest date>"]
    for name in ["Base", "Promo", "Media", "Blanks", "Positives", "Negatives"]:
        chart_data.add_series(name, (0, 0, 0))
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(3),
        chart_data,
    )
    return chart_shape.chart

def deck_automation(request):
    context: dict[str, object] = {
        "template_options": sorted(TEMPLATE_OPTIONS.keys()),
        "selected_template": "MMx",
        "year1_value": "Year1",
        "year2_value": "Year2",
        "bucket_config_json": "",
    }
    if request.method == "POST":
        gathered_file = request.FILES.get("gathered_cn10")
        scope_file = request.FILES.get("scope_workbook")
        selected_template = request.POST.get("template_choice", "").strip() or "MMx"
        context["selected_template"] = selected_template

        year1_value = request.POST.get("year1", "Year1").strip() or "Year1"
        year2_value = request.POST.get("year2", "Year2").strip() or "Year2"
        bucket_config_json = request.POST.get("bucket_config_json", "").strip()
        context["year1_value"] = year1_value
        context["year2_value"] = year2_value
        context["bucket_config_json"] = bucket_config_json

        template_filename = TEMPLATE_OPTIONS.get(selected_template, "MMx.pptx")

        if not gathered_file:
            context["error"] = "Please upload the gatheredCN10 file to continue."
            return render(request, "deck_automation/deck_automation.html", context)

        try:
            template_path = Path(__file__).resolve().parent / "templates" / "deck_automation" / template_filename
            if not template_path.exists():
                raise FileNotFoundError(f"Template not found: {template_filename}")

            gathered_df = read_df(gathered_file)
            scope_df = read_scope_df(scope_file)
            template_chart = _find_template_chart(str(template_path)) or _default_template_chart()

            bucket_data = None
            if bucket_config_json:
                bucket_config = json.loads(bucket_config_json)
                if not isinstance(bucket_config, dict):
                    raise ValueError("Bucket config must be a JSON object keyed by bucket label.")
                bucket_data = {
                    "year1": year1_value,
                    "year2": year2_value,
                    "bucket_config": bucket_config,
                }

            payloads_by_label = compute_waterfall_payloads_for_all_labels(
                gathered_df,
                scope_df,
                bucket_data=bucket_data,
                template_chart=template_chart,
            )
            payload_json = waterfall_payloads_to_json(payloads_by_label)

            download_id = str(uuid4())
            payloads = request.session.get(SESSION_PAYLOADS_KEY, {})
            payloads[download_id] = payload_json
            request.session[SESSION_PAYLOADS_KEY] = payloads
            request.session.modified = True

            summaries = []
            for label, payload in payloads_by_label.items():
                summaries.append(
                    {
                        "label": label,
                        "category_count": len(payload.categories),
                        "checksum": payload_checksum(payload),
                    }
                )


            if settings.DEBUG:
                first_label = next(iter(payloads_by_label), None)
                first_payload_categories = payloads_by_label[first_label].categories[:5] if first_label else []
                context["debug_gathered_header_row_index"] = gathered_df.attrs.get("detected_header_row_index")
                context["debug_gathered_columns"] = list(gathered_df.columns)
                context["debug_first_payload_categories"] = first_payload_categories

            context.update(
                {
                    "computed_count": len(payloads_by_label),
                    "payload_summaries": summaries,
                    "download_id": download_id,
                    "message": "Payload computation complete.",
                    "selected_template": selected_template,
                }
            )
        except Exception as exc:  # noqa: BLE001
            logger.exception("Deck automation payload computation failed")
            context["error"] = str(exc) or "We could not compute payloads from the uploaded file."

    return render(request, "deck_automation/deck_automation.html", context)


def download_payloads_json(request, download_id: str):
    payloads = request.session.get(SESSION_PAYLOADS_KEY, {})
    payload_json = payloads.get(str(download_id))
    if not payload_json:
        raise Http404("Payloads not found.")
    response = HttpResponse(payload_json, content_type="application/json")
    response["Content-Disposition"] = "attachment; filename=payloads.json"
    return response
